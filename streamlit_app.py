import streamlit as st
import pandas as pd
import base64
import os
import io
import re
from datetime import datetime, date, timedelta
import pytz

st.set_page_config(page_title="AMC NTEP Dashboard", layout="wide", initial_sidebar_state="collapsed")

def img_to_b64(img_path):
    try:
        with open(img_path, "rb") as img_file: return base64.b64encode(img_file.read()).decode('utf-8')
    except: return ""

LOG_FILE = "activity_log.csv"
india_tz = pytz.timezone('Asia/Kolkata')

def log_activity(username, role, target, action):
    if not os.path.exists(LOG_FILE):
        df_log = pd.DataFrame(columns=["Timestamp", "Username", "Role", "Target", "Action"])
        df_log.to_csv(LOG_FILE, index=False)
    current_time = datetime.now(india_tz).strftime("%d-%b-%Y, %I:%M %p")
    new_entry = pd.DataFrame([{"Timestamp": current_time, "Username": username, "Role": role, "Target": target, "Action": action}])
    new_entry.to_csv(LOG_FILE, mode='a', header=False, index=False)

if "auth" not in st.session_state: 
    st.session_state.auth = False
    st.session_state.current_user = ""
    st.session_state.role = ""
    st.session_state.target = ""

try:
    df_users = pd.read_csv("users.csv")
    df_users['Username'] = df_users['Username'].astype(str).str.strip().str.upper()
    df_users['Password'] = df_users['Password'].astype(str).str.strip()
except:
    st.error("⚠️ User Database (users.csv) મળ્યું નથી!")
    st.stop()

# ==========================================
# 🔐 NEW ENTERPRISE LOGIN PAGE DESIGN (BUG FIXED)
# ==========================================
if not st.session_state.auth:
    b64_amc = img_to_b64("images/amc.png")
    
    st.markdown("""
    <style>
    .left-panel { background: #0A3A6E; color: white; padding: 40px 30px; border-radius: 15px 0 0 15px; height: 100%; text-align: center; position: relative; overflow: hidden; }
    .right-panel { padding: 40px; background: white; border-radius: 0 15px 15px 0; border: 1px solid #e2e8f0; border-left: none; height: 100%; display: flex; flex-direction: column; justify-content: center;}
    .stTextInput>div>div>input { background-color: #f8fafc; border-radius: 8px; border: 1px solid #cbd5e1; padding: 12px; }
    .stButton>button { background-color: #0A3A6E; color: white; border-radius: 8px; width: 100%; font-weight: 600; padding: 10px; margin-top: 15px; }
    .stButton>button:hover { background-color: #185FA5; color: white; border-color: #185FA5; }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown("<br><br>", unsafe_allow_html=True)
    sp1, login_box, sp2 = st.columns([1, 6, 1])
    
    with login_box:
        l_col, r_col = st.columns([4, 5], gap="small")
        
        with l_col:
            st.markdown(f"""
<div class="left-panel">
<div style="background: #F0F6FF; width: 85px; height: 85px; border-radius: 50%; margin: 0 auto 20px auto; display: flex; align-items: center; justify-content: center; border: 2px solid rgba(255,255,255,0.25);">
<img src="data:image/png;base64,{b64_amc}" width="65">
</div>
<div style="display:inline-block; background:rgba(255,255,255,0.1); padding: 5px 13px; border-radius: 20px; font-size: 10px; letter-spacing: 1.8px; margin-bottom: 20px; border: 0.5px solid rgba(255,255,255,0.2);">
<span style="color:#5DCAA5;">●</span> <span style="color:#9FC8F0; font-weight:600;">AMC · NTEP</span>
</div>
<h2 style="font-family: serif; margin-bottom: 12px; font-size: 24px; font-weight: 600;">National TB Elimination Programme</h2>
<p style="font-size: 13px; color: #85B7EB; line-height: 1.6; margin-bottom:30px;">Ahmedabad Municipal Corporation's centralised surveillance & management platform for TB programme monitoring.</p>
<div style="border-top: 1px solid rgba(255,255,255,0.14); padding-top: 15px; display: flex; justify-content: space-around;">
<div><b style="font-size: 18px;">7</b><br><span style="font-size: 10px; color: #85B7EB;">ZONES</span></div>
<div><b style="font-size: 18px;">23</b><br><span style="font-size: 10px; color: #85B7EB;">TB UNITS</span></div>
<div><b style="font-size: 18px;">Live</b><br><span style="font-size: 10px; color: #85B7EB;">REPORTING</span></div>
</div>
</div>
            """, unsafe_allow_html=True)
            
        with r_col:
            st.markdown("""
<div style="padding: 10px 10px 20px 10px;">
<h3 style="color: #1e293b; margin-bottom: 5px; font-weight: 600;">Sign in to your account</h3>
<p style="color: #64748b; font-size: 13px;">Access restricted to authorised Zone & TB Unit personnel only.</p>
</div>
            """, unsafe_allow_html=True)
            
            uname = st.text_input("User ID / Zone Code", placeholder="e.g. AMC-Z3-001").strip().upper()
            pwd = st.text_input("Password", type="password", placeholder="Enter your password").strip()
            
            if st.button("Sign In Securely", use_container_width=True):
                user_match = df_users[(df_users['Username'] == uname) & (df_users['Password'] == pwd)]
                if not user_match.empty: 
                    st.session_state.auth = True
                    st.session_state.current_user = uname
                    st.session_state.role = user_match.iloc[0]['Role']
                    st.session_state.target = user_match.iloc[0]['Target']
                    log_activity(st.session_state.current_user, st.session_state.role, st.session_state.target, "Logged In")
                    st.rerun()
                else: 
                    st.error("⚠️ Invalid User ID or Password")
                    
            st.markdown("<p style='text-align: right; color: #378ADD; font-size: 12px; margin-top: 15px; cursor: pointer;'>Forgot password?</p>", unsafe_allow_html=True)

    st.markdown("<p style='text-align: center; color: #94a3b8; font-size: 12px; margin-top: 25px;'>© 2026 Ahmedabad Municipal Corporation · All rights reserved</p>", unsafe_allow_html=True)
    st.stop()

# ==========================================
# 🟢 POST LOGIN DASHBOARD UI & DATA FETCHING (તમારો જૂનો ડેટા કોડ)
# ==========================================
st.markdown("""<style>#MainMenu {visibility: hidden;} header {visibility: hidden;} footer {visibility: hidden;}</style>""", unsafe_allow_html=True)
st.markdown(f"<div style='background-color: #d4edda; color: #155724; padding: 12px; border-radius: 8px; border: 1px solid #c3e6cb; margin-bottom: 10px; font-size: 16px; font-weight: bold;'>👤 Logged in as: <span style='color: #0b2e13; font-size: 18px;'>{st.session_state.target} ({st.session_state.role})</span></div>", unsafe_allow_html=True)

with st.expander("⚙️ Account Settings & Change Password"):
    c_p0, c_p1, c_p2, c_p3 = st.columns([2, 2, 2, 1])
    with c_p0: old_pwd = st.text_input("Old Password", type="password", key="p0")
    with c_p1: new_pwd = st.text_input("New Password", type="password", key="p1")
    with c_p2: conf_pwd = st.text_input("Confirm Password", type="password", key="p2")
    with c_p3:
        st.write(""); st.write("")
        if st.button("Update", use_container_width=True):
            current_actual_pwd = df_users.loc[df_users['Username'] == st.session_state.current_user, 'Password'].values[0]
            if old_pwd != current_actual_pwd: st.error("⚠️ Old Password is incorrect!")
            elif new_pwd != conf_pwd: st.error("⚠️ New Passwords do not match!")
            elif new_pwd == "": st.error("⚠️ Password cannot be empty!")
            else:
                df_users.loc[df_users['Username'] == st.session_state.current_user, 'Password'] = new_pwd
                df_users.to_csv("users.csv", index=False)
                st.success("✅ Password updated!")

if st.session_state.role == "ADMIN":
    with st.expander("🛡️ Admin Panel: View Passwords & Activity Logs"):
        a_tab1, a_tab2 = st.tabs(["🔑 Manage Users", "📝 Activity Logs"])
        with a_tab1:
            st.dataframe(df_users, use_container_width=True, hide_index=True)
        with a_tab2:
            try:
                df_logs = pd.read_csv(LOG_FILE)
                st.dataframe(df_logs.iloc[::-1], use_container_width=True, hide_index=True)
            except: st.write("No logs available yet.")

st.markdown("<hr style='margin-top: 5px; margin-bottom: 15px;'>", unsafe_allow_html=True)

# 🎯 તમારું EXCEL ડાઉનલોડ કરવાનું ફંક્શન
def convert_df_to_excel(df, sheet_name="Data"):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name=sheet_name)
        workbook = writer.book
        worksheet = writer.sheets[sheet_name]
        header_format = workbook.add_format({'bold': True, 'text_wrap': True, 'valign': 'top', 'align': 'center', 'fg_color': '#1f618d', 'font_color': 'white', 'border': 1})
        cell_format = workbook.add_format({'align': 'center', 'valign': 'vcenter', 'border': 1})
        for col_num, value in enumerate(df.columns.values):
            worksheet.write(0, col_num, value, header_format)
        for i, col in enumerate(df.columns):
            if len(df) == 0: column_len = len(str(col)) + 2
            else:
                max_val_len = df[col].astype(str).str.len().max()
                column_len = max(max_val_len if pd.notna(max_val_len) else 0, len(str(col))) + 2
            if column_len > 30: column_len = 30 
            worksheet.set_column(i, i, int(column_len), cell_format)
    return output.getvalue()

# 🎯 તમારું DATA FETCH કરવાનું ફંક્શન
@st.cache_data(ttl=3600)
def load_all_data():
    try:
        m = pd.read_csv("Master_Line_List.csv", dtype={'Episode ID': str})
        for c in ['Diagnosis Date', 'Initiation Date', 'Outcome Date']:
            if c in m.columns: m[c] = pd.to_datetime(m[c], errors='coerce') 
        c_mat = pd.read_csv("Comparison_Matrix.csv", dtype={'Episode ID': str})
        if not c_mat.empty and not m.empty:
            dates_df = m[['Episode ID', 'Diagnosis Date', 'Initiation Date', 'Outcome Date']].drop_duplicates('Episode ID')
            c_mat = c_mat.merge(dates_df, on='Episode ID', how='left')
        curr = pd.read_csv("Current_TB_Patients.csv", dtype={'Episode ID': str})
        t_df = pd.read_csv("Update_Timestamps.csv")
        try:
            p_today = pd.read_csv("Presumptive_Today.csv", dtype={'Episode_ID': str})
            p_yest = pd.read_csv("Presumptive_Yest.csv", dtype={'Episode_ID': str})
        except:
            p_today, p_yest = pd.DataFrame(), pd.DataFrame()
            
        return m, c_mat, curr, t_df, p_today, p_yest
    except: return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

# 🎯 તમારું ગૂગલ શીટ માંથી LIVE DATA લાવવાનું ફંક્શન
@st.cache_data(ttl=300) 
def get_live_dc():
    def fetch_sheet(url):
        try:
            df = pd.read_csv(url, header=None, low_memory=False, dtype=str)
            header_row = -1
            for i in range(min(20, len(df))):
                row_str = " ".join(df.iloc[i].fillna("").astype(str).str.upper())
                if "EPISODE" in row_str and "NAME" in row_str:
                    header_row = i; break
            
            if header_row == -1: return pd.DataFrame()
                
            header_vals = df.iloc[header_row].fillna("").astype(str).str.upper()
            df = df.iloc[header_row+1:].reset_index(drop=True)
            
            def cx_col(col_let):
                num = 0
                for c in col_let.upper(): num = num * 26 + (ord(c) - ord('A') + 1)
                return num - 1

            tu_idx, phi_idx, id_idx, name_idx, zone_idx = cx_col('A'), cx_col('C'), cx_col('G'), cx_col('H'), cx_col('AR')
            diag_idx, init_idx, out_idx = cx_col('CJ'), cx_col('CK'), cx_col('CL')
            hf_idx, case_idx, site_idx, out_col_idx = cx_col('B'), cx_col('Z'), cx_col('AA'), cx_col('AD')
            ci_idx = cx_col('CI') 
            cx_base, cy_1m, cz_2m, da_3m, db_4m, dc_5m, dd_6m = cx_col('CX'), cx_col('CY'), cx_col('CZ'), cx_col('DA'), cx_col('DB'), cx_col('DC'), cx_col('DD')

            for i, val in enumerate(header_vals):
                val_c = val.strip()
                if "EPISODE" in val_c and "ID" in val_c: id_idx = i
                elif "PATIENT" in val_c and "NAME" in val_c: name_idx = i
                elif "DUE" in val_c and "STATUS" in val_c: ci_idx = i
                elif "DIAGNOSIS" in val_c and "DATE" in val_c: diag_idx = i
                elif "INITIATION" in val_c and "DATE" in val_c: init_idx = i
                elif "OUTCOME" in val_c and "DATE" in val_c: out_idx = i
                elif val_c == "ZONE": zone_idx = i
                elif "ELIGIBILITY" in val_c and "BASE" in val_c: cx_base = i
                elif "ELIGIBILITY" in val_c and "1" in val_c: cy_1m = i
                elif "ELIGIBILITY" in val_c and "2" in val_c: cz_2m = i
                elif "ELIGIBILITY" in val_c and "3" in val_c: da_3m = i
                elif "ELIGIBILITY" in val_c and "4" in val_c: db_4m = i
                elif "ELIGIBILITY" in val_c and "5" in val_c: dc_5m = i
                elif "ELIGIBILITY" in val_c and "6" in val_c: dd_6m = i

            diff_data = []
            for _, row in df.iterrows():
                def get_v(idx): return str(row.iloc[idx]).strip().upper() if idx < len(row) else ""
                elig_base, elig_1m, elig_2m, elig_3m, elig_4m, elig_5m, elig_6m = get_v(cx_base), get_v(cy_1m), get_v(cz_2m), get_v(da_3m), get_v(db_4m), get_v(dc_5m), get_v(dd_6m)
                is_elig = any("ELIG" in v and "NOT" not in v for v in [elig_base, elig_1m, elig_2m, elig_3m, elig_4m, elig_5m, elig_6m])
                if is_elig:
                    tu = get_v(tu_idx).replace("-", "")
                    if "INDIA" in tu: tu = "INDIA COLONY"
                    elif "NAVA" in tu and "VADAJ" in tu: tu = "NAVA VADAJ"
                    elif "JUNA" in tu and "VADAJ" in tu: tu = "JUNA VADAJ"
                    elif "NOB" in tu: tu = "NOBLENAGAR"
                    elif "BEHRAM" in tu: tu = "BEHRAMPURA"
                    elif "SAIJ" in tu: tu = "SAIJPUR"
                    elif "DANI" in tu: tu = "DANILIMDA"
                    elif "AMRAI" in tu: tu = "AMRAIWADI"
                    elif "BHAI" in tu: tu = "BHAIPURA"
                    elif "GHAT" in tu: tu = "GHATLODIA"
                    elif "CHAND" in tu: tu = "CHANDKHEDA"
                    elif "VEJAL" in tu: tu = "VEJALPUR"
                    elif "ISAN" in tu: tu = "ISANPUR"
                    elif "ASAR" in tu: tu = "ASARVA"
                    elif "BAPU" in tu: tu = "BAPUNAGAR"
                    elif "VIRAT" in tu: tu = "VIRATNAGAR"
                    elif "RAKH" in tu: tu = "RAKHIAL"
                    elif "JAMAL" in tu: tu = "JAMALPUR"
                    elif "VASNA" in tu: tu = "VASNA"
                    elif "VATVA" in tu: tu = "VATVA"
                    elif "JODH" in tu: tu = "JODHPUR"
                    elif "SHAH" in tu: tu = "SHAHPUR"
                    elif "RANIP" in tu: tu = "RANIP"
                    zone = get_v(zone_idx)
                    if zone in ["", "NAN", "NONE", "NULL", "N/A"]: zone = 'MAPPING NOT DONE'
                    diff_data.append({
                        'ZONE': zone, 'TB Unit': tu, 'PHI': get_v(phi_idx), 'Episode ID': get_v(id_idx), 'Patient Name': get_v(name_idx),
                        'Due_Status': get_v(ci_idx), 'Diagnosis Date': get_v(diag_idx), 'Initiation Date': get_v(init_idx), 'Outcome Date': get_v(out_idx),
                        'Facility_Type': get_v(hf_idx), 'Type_of_Case': get_v(case_idx), 
                        'Site_of_TBDisease': get_v(site_idx), 'Treatment_Outcome': get_v(out_col_idx),
                        'Elig_BASELINE': elig_base, 'Elig_1ST_MONTH': elig_1m, 'Elig_2ND_MONTH': elig_2m,
                        'Elig_3RD_MONTH': elig_3m, 'Elig_4TH_MONTH': elig_4m, 'Elig_5TH_MONTH': elig_5m, 'Elig_6TH_MONTH': elig_6m
                    })
            return pd.DataFrame(diff_data)
        except Exception as e:
            return pd.DataFrame()

    try:
        url_new = "https://docs.google.com/spreadsheets/d/1hkJBnJOuxcVu233f6e2_0cOE-BM7bdDOyHuzrlGogMU/export?format=csv&gid=1152778583"
        url_old = "https://docs.google.com/spreadsheets/d/1zdf96eisZHzdk5ECFSI7eeOtNQoOXk3QRUUROtIZQmc/export?format=csv&gid=1152778583"
        
        df_new = fetch_sheet(url_new)
        df_old = fetch_sheet(url_old)
        
        return df_new, df_old
    except:
        return pd.DataFrame(), pd.DataFrame()

df_master_raw, df_comp_raw, df_curr_tb_raw, df_time, df_pres_t_raw, df_pres_y_raw = load_all_data()
df_dc_new_raw, df_dc_old_raw = get_live_dc()

# ==========================================
# 🎯 BULLETPROOF LOGIN FILTER 
# ==========================================
def filter_by_role(df, role, target):
    if df.empty: return df
    target_up = str(target).upper().strip()
    role_up = str(role).upper().strip()
    
    if role_up in ["TB_UNIT", "TB UNIT", "TU"]:
        tu_col = 'TB Unit' if 'TB Unit' in df.columns else 'TB_UNIT' if 'TB_UNIT' in df.columns else None
        if tu_col:
            def strict_tu_check(val):
                v = str(val).upper().strip()
                if target_up == "VADAJ" and ("JUNA" in v or "NAVA" in v): return False
                if target_up == "RANIP" and "NEW" in v: return False
                return target_up in v
            return df[df[tu_col].apply(strict_tu_check)]
            
    elif role_up == "ZONE" and 'ZONE' in df.columns:
        target_clean = target_up.replace("ZONE", "").strip()
        def strict_zone_check(val):
            v_raw = str(val).upper().strip()
            v_list = [z.strip().replace("ZONE", "").strip() for z in v_raw.replace(',', '&').split('&')]
            return target_clean in v_list
        return df[df['ZONE'].apply(strict_zone_check)]
        
    return df

df_master = filter_by_role(df_master_raw.copy(), st.session_state.role, st.session_state.target)
df_comp = filter_by_role(df_comp_raw.copy(), st.session_state.role, st.session_state.target) 
df_curr_tb = filter_by_role(df_curr_tb_raw.copy(), st.session_state.role, st.session_state.target)
df_dc_new = filter_by_role(df_dc_new_raw.copy(), st.session_state.role, st.session_state.target)
df_dc_old = filter_by_role(df_dc_old_raw.copy(), st.session_state.role, st.session_state.target)
df_pres_t = filter_by_role(df_pres_t_raw.copy(), st.session_state.role, st.session_state.target)
df_pres_y = filter_by_role(df_pres_y_raw.copy(), st.session_state.role, st.session_state.target)

def draw_card(title, value, color, icon):
    return f"""<div style="background-color: {color}; border-radius: 8px; padding: 15px 5px; margin-bottom: 10px; color: white; text-align: center; box-shadow: 0 4px 6px rgba(0,0,0,0.1);"><div style="font-size: 24px; margin-bottom: 5px;">{icon}</div><div style="font-size: 13px; font-weight: bold; text-transform: uppercase;">{title}</div><div style="font-size: 26px; font-weight: 900; margin-top: 8px;">{value}</div></div>"""

def clean_selection(selected_list): return [item.rsplit(" (", 1)[0] for item in selected_list]

def get_options_with_counts(df, column_name, tab_name="tab1"):
    if df.empty or column_name not in df.columns: return []
    try:
        if tab_name == "tab1" and 'Pending Status' in df.columns:
            df_temp = df.copy()
            df_temp['act_cnt'] = df_temp['Pending Status'].astype(str).apply(lambda x: len([s for s in x.split('+') if s.strip()]))
            counts = df_temp.groupby(column_name)['act_cnt'].sum()
        else:
            counts = df[column_name].value_counts()
        counts = counts[counts > 0].sort_values(ascending=False)
        return [f"{val} ({int(count)})" for val, count in counts.items() if str(val) not in ["nan", "", "None", "N/A"]]
    except: return []

b64_amc, b64_ntep = img_to_b64("images/amc.png"), img_to_b64("images/ntep.jpg")
st.markdown(f"<div style='display: flex; justify-content: space-between; align-items: center;'><img src='data:image/png;base64,{b64_amc}' height='75'><h3 style='margin:0; font-weight:900;'>AMC | NTEP</h3><img src='data:image/jpeg;base64,{b64_ntep}' height='75'></div>", unsafe_allow_html=True)
st.markdown("<div style='background-color:#1f618d; color:white; text-align:center; padding:12px; border-radius:5px; margin:15px 0;'>TB Monitoring Dashboard - Ahmedabad</div>", unsafe_allow_html=True)

if not df_time.empty:
    with st.expander("🕒 Register Last Sync Timestamps (IST)"):
        t_cols = st.columns(6)
        for i, row in df_time.iterrows():
            color = "#27AE60" if "Live" in str(row['Last Updated']) else "#E67E22"
            with t_cols[i % 6]: 
                st.markdown(f"<div style='font-size:13px; color:#333;'><b>{row['Register']}</b><br><span style='color:{color}; font-weight:bold;'>{row['Last Updated']}</span></div>", unsafe_allow_html=True)

tab1, tab2, tab4, tab5, tab6, tab7, tab8, tab9 = st.tabs(["📊 Master Dashboard", "🔄 Daily Comparison", "🚀 Smart PPT", "🏥 Diff. Care", "👥 Staff Directory", "🔬 Presumptive TB", "🚨 Adverse Outcomes", "📱 Live Field Data"])


# ==========================================
# 🟢 TAB 1: MASTER DASHBOARD
# ==========================================
with tab1:
    with st.expander("🔽 Filters & Sorting"):
        c1, c2, c3 = st.columns(3)
        df_disp = df_master.copy()
        with c1:
            if st.session_state.role == "ADMIN":
                s_z = clean_selection(st.multiselect("Zone", get_options_with_counts(df_disp, 'ZONE', 'tab1'), key='z1'))
                if s_z: df_disp = df_disp[df_disp['ZONE'].isin(s_z)]
            # 🎯 DEPENDENT FILTER: TB Unit based on Zone
            s_tu = clean_selection(st.multiselect("TB Unit", get_options_with_counts(df_disp, 'TB Unit', 'tab1'), key='tu1'))
            if s_tu: df_disp = df_disp[df_disp['TB Unit'].isin(s_tu)]
        with c2:
            if 'Facility Type' in df_disp.columns:
                available_facs = df_disp['Facility Type'].astype(str).str.upper().unique()
                fac_opts = [f for f in ["PUBLIC", "PRIVATE"] if any(a in ["PUBLIC", "PHI"] if f=="PUBLIC" else a not in ["PUBLIC", "PHI", "N/A", "NAN", ""] for a in available_facs)]
                s_ft_raw = st.multiselect("Facility Category", fac_opts, key='fc1')
                if s_ft_raw:
                    if "PUBLIC" in s_ft_raw and "PRIVATE" in s_ft_raw: pass
                    elif "PUBLIC" in s_ft_raw: df_disp = df_disp[df_disp['Facility Type'].astype(str).str.upper().isin(['PUBLIC', 'PHI'])]
                    elif "PRIVATE" in s_ft_raw: df_disp = df_disp[~df_disp['Facility Type'].astype(str).str.upper().isin(['PUBLIC', 'PHI'])]
            # 🎯 DEPENDENT FILTER: PHI based on TB Unit
            s_phi = clean_selection(st.multiselect("Filter PHI", get_options_with_counts(df_disp, 'PHI', 'tab1'), key='phi1'))
            if s_phi: df_disp = df_disp[df_disp['PHI'].isin(s_phi)]
            
            inds = ["Outcome", "UDST", "Not Put On", "SLPA", "Consent", "ADT", "RBS", "ART", "CPT", "HIV"]
            f_rep = st.multiselect("Report Type", inds, key='rep1')
        with c3:
            diag_dt = st.date_input("Diagnosis Date Range", value=[], key="d1")
            init_dt = st.date_input("Initiation Date Range", value=[], key="d2")
            out_dt = st.date_input("Outcome Date Range", value=[], key="d3")
            
        if len(diag_dt) == 2: df_disp = df_disp[pd.to_datetime(df_disp.get('Diagnosis Date'), errors='coerce').dt.date.between(diag_dt[0], diag_dt[1])]
        if len(init_dt) == 2: df_disp = df_disp[pd.to_datetime(df_disp.get('Initiation Date'), errors='coerce').dt.date.between(init_dt[0], init_dt[1])]
        if len(out_dt) == 2: df_disp = df_disp[pd.to_datetime(df_disp.get('Outcome Date'), errors='coerce').dt.date.between(out_dt[0], out_dt[1])]
        if f_rep and 'Pending Status' in df_disp.columns: df_disp = df_disp[df_disp['Pending Status'].str.contains("|".join(f_rep), na=False)]

    if 'Pending Status' in df_disp.columns:
        f_counts = {k: len(df_disp[df_disp['Pending Status'].str.contains(k, na=False)]) for k in inds}
        sorted_counts = sorted(f_counts.items(), key=lambda x: x[1], reverse=True)
        top_3, others = sorted_counts[:3], sorted_counts[3:]
        colors = {"Outcome": "#F39C12", "UDST": "#C0392B", "Not Put On": "#27AE60", "SLPA": "#8E44AD", "Consent": "#D35400", "HIV": "#C0392B", "ART": "#2980B9", "CPT": "#2980B9", "RBS": "#16A085", "ADT": "#E67E22"}
        st.markdown("##### 📈 Top 3 Highest Pending Actions")
        cc1, cc2, cc3, cc4 = st.columns(4)
        with cc1: st.markdown(draw_card("Total Pendency", sum(f_counts.values()), "#1f618d", "📄"), unsafe_allow_html=True)
        with cc2: st.markdown(draw_card(top_3[0][0], top_3[0][1], colors.get(top_3[0][0], "#34495E"), "📌"), unsafe_allow_html=True)
        with cc3: st.markdown(draw_card(top_3[1][0], top_3[1][1], colors.get(top_3[1][0], "#34495E"), "📌"), unsafe_allow_html=True)
        with cc4: st.markdown(draw_card(top_3[2][0], top_3[2][1], colors.get(top_3[2][0], "#34495E"), "📌"), unsafe_allow_html=True)
        with st.expander("🔽 Tap to show other reports"):
            oc_cols = st.columns(4)
            for i, (k, v) in enumerate(others):
                with oc_cols[i % 4]: st.markdown(draw_card(k, v, colors.get(k, "#34495E"), "📌"), unsafe_allow_html=True)
    
    # 👇👇👇 CLINICAL STATUS ENGINE 👇👇👇
    if not df_disp.empty:
        with st.spinner("Calculating Clinical Status..."):
            import re
            
            # 1. VISUAL CLEANER: Wipe fake string values safely
            for c in ['Diagnosis Date', 'Initiation Date', 'Outcome Date', 'Treatment Outcome', 'Pending Status', 'Extend Status']:
                if c in df_disp.columns:
                    df_disp[c] = df_disp[c].astype(str).str.strip()
                    df_disp[c] = df_disp[c].replace(['None', 'nan', 'NaN', 'N/A', '<NA>', 'NaT', ''], pd.NA)

            def get_dt(col_name):
                if col_name in df_disp.columns:
                    return pd.to_datetime(df_disp[col_name], errors='coerce')
                return pd.Series([pd.NaT]*len(df_disp), index=df_disp.index)

            diag_dt = get_dt('Diagnosis Date')
            init_dt = get_dt('Initiation Date')
            out_dt = get_dt('Outcome Date')
            today = pd.Timestamp.today().normalize()

            if 'Treatment Outcome' in df_disp.columns:
                out_str = df_disp['Treatment Outcome'].fillna("").astype(str).str.upper()
                has_outcome = ~out_str.isin(["", "NAN", "NAT", "NONE", "NULL", "<NA>", "N/A"])
            else:
                has_outcome = pd.Series([False]*len(df_disp), index=df_disp.index)

            # Initialize empty columns (Uses blank text to prevent Pandas Crash)
            df_disp['Treatment Status'] = pd.Series("", index=df_disp.index)
            df_disp['On Treatment Days'] = pd.Series("", index=df_disp.index)

            # 🎯 MASK: Only apply math if they have valid notification dates OR an outcome OR a notification-related pending status.
            has_notif_data = diag_dt.notna() | init_dt.notna() | has_outcome | df_disp['Pending Status'].fillna("").astype(str).str.upper().str.contains("OUTCOME|NOT PUT ON", na=False)

            # --- 🎯 APPLY YOUR EXACT CLINICAL RULES ---

            # 🔴 RULE 1: Not Put On (Diag is present, Initiation is Blank, Outcome is Blank)
            mask_npo = has_notif_data & diag_dt.notna() & init_dt.isna() & (~has_outcome)
            df_disp.loc[mask_npo, 'Treatment Status'] = "Not Put On"
            # Colab puts "Not Put On" in Pending Status natively, but we enforce it here visually just in case
            if 'Pending Status' in df_disp.columns:
                curr_pend = df_disp.loc[mask_npo, 'Pending Status'].astype(str)
                needs_add = ~curr_pend.str.contains("Not Put On", case=False, na=False)
                df_disp.loc[mask_npo & needs_add, 'Pending Status'] = df_disp.loc[mask_npo & needs_add, 'Pending Status'].apply(lambda x: f"{x} + Not Put On" if str(x).strip() else "Not Put On")
                df_disp['Pending Status'] = df_disp['Pending Status'].str.replace(r'^\+\s*', '', regex=True)

            # 🔴 RULE 2: Initial Defaulter (Diag Exists, Init Blank, AND OUTCOME IS PRESENT)
            mask_defaulter = has_notif_data & diag_dt.notna() & init_dt.isna() & has_outcome
            df_disp.loc[mask_defaulter, 'Treatment Status'] = "Initial Defaulter"

            # 🟢 RULE 3: Treatment Given (Initiation Date Exists)
            mask_given = has_notif_data & init_dt.notna()
            df_disp.loc[mask_given, 'Treatment Status'] = "Treatment Given"

            # ⏳ RULE 4: On Treatment Days (Condition 1 - Outcome is BLANK -> Count to Today)
            mask_days_blank_out = mask_given & (~has_outcome)
            if mask_days_blank_out.any():
                df_disp.loc[mask_days_blank_out, 'On Treatment Days'] = (today - init_dt[mask_days_blank_out]).dt.days.astype(int).astype(str) + " Days"

            # ⏳ RULE 5: On Treatment Days (Condition 2 - Outcome is PRESENT -> Count to Outcome Date)
            mask_days_has_out = mask_given & has_outcome & out_dt.notna()
            if mask_days_has_out.any():
                df_disp.loc[mask_days_has_out, 'On Treatment Days'] = (out_dt[mask_days_has_out] - init_dt[mask_days_has_out]).dt.days.astype(int).astype(str) + " Days"

            # 🧹 CLEANUP: If Initial Defaulter, ERASE "Not Put On" from their Pending Status
            if 'Pending Status' in df_disp.columns:
                cleaned_pend = df_disp.loc[mask_defaulter, 'Pending Status'].astype(str).str.replace('Not Put On', '', flags=re.IGNORECASE)
                cleaned_pend = cleaned_pend.str.replace(r'\+\s*\+', '+', regex=True).str.strip(' +')
                df_disp.loc[mask_defaulter, 'Pending Status'] = cleaned_pend

            # 2. Final Visual Polish
            df_disp = df_disp.fillna('')
            if 'Pending Status' in df_disp.columns: df_disp['Pending Status'] = df_disp['Pending Status'].replace(['<NA>', 'nan', 'None'], '')
            
            # 🔥 Ensure "Extend Status" displays properly
            if 'Extend Status' in df_disp.columns:
                df_disp['Extend Status'] = df_disp['Extend Status'].astype(str).replace(['<NA>', 'nan', 'None', 'N/A', 'NAN', ''], '')

            # 3. Insert Columns perfectly next to Treatment Outcome
            cols = df_disp.columns.tolist()
            if 'Treatment Status' in cols: cols.remove('Treatment Status')
            if 'On Treatment Days' in cols: cols.remove('On Treatment Days')
            if 'Treatment Outcome' in cols:
                insert_idx = cols.index('Treatment Outcome') + 1
                cols.insert(insert_idx, 'Treatment Status')
                cols.insert(insert_idx + 1, 'On Treatment Days')
                df_disp = df_disp[cols]
    # 👆👆👆 END PERFECT CLINICAL STATUS ENGINE 👆👆👆

    st.markdown(f"<div style='color: #2E86C1; margin-bottom: 10px; font-weight: bold;'>Found {len(df_disp)} Patient(s)</div>", unsafe_allow_html=True)
    st.dataframe(df_disp, use_container_width=True, hide_index=True)
    if not df_disp.empty:
        st.download_button("📥 Download Master Excel", convert_df_to_excel(df_disp, "Master_Report"), "Master_Report.xlsx", key='dl1')                

# ==========================================
# 🟢 TAB 2: DAILY COMPARISON (NO DATES DISPLAYED)
# ==========================================
with tab2:
    st.markdown("#### 🔄 Comparison Matrix")
    with st.expander("🔽 Filters & Dates", expanded=True):
        c1, c2, c3 = st.columns(3)
        df_c = df_comp.copy()
        with c1: 
            if st.session_state.role == "ADMIN":
                s2_z = clean_selection(st.multiselect("Filter Zone", get_options_with_counts(df_c, 'ZONE', 'tab2'), key='z2'))
                if s2_z: df_c = df_c[df_c['ZONE'].isin(s2_z)]
            # 🎯 DEPENDENT FILTER
            s2_tu = clean_selection(st.multiselect("Filter TB Unit", get_options_with_counts(df_c, 'TB Unit', 'tab2'), key='tu2'))
            if s2_tu: df_c = df_c[df_c['TB Unit'].isin(s2_tu)]
        with c2: 
            if 'Facility Type' in df_c.columns:
                available_facs2 = df_c['Facility Type'].astype(str).str.upper().unique()
                fac_opts2 = [f for f in ["PUBLIC", "PRIVATE"] if any(a in ["PUBLIC", "PHI"] if f=="PUBLIC" else a not in ["PUBLIC", "PHI", "N/A", "NAN", ""] for a in available_facs2)]
                s2_ft_raw = st.multiselect("Facility Category", fac_opts2, key='fc2')
                if s2_ft_raw:
                    if "PUBLIC" in s2_ft_raw and "PRIVATE" in s2_ft_raw: pass
                    elif "PUBLIC" in s2_ft_raw: df_c = df_c[df_c['Facility Type'].astype(str).str.upper().isin(['PUBLIC', 'PHI'])]
                    elif "PRIVATE" in s2_ft_raw: df_c = df_c[~df_c['Facility Type'].astype(str).str.upper().isin(['PUBLIC', 'PHI'])]
            # 🎯 DEPENDENT FILTER
            s2_phi = clean_selection(st.multiselect("Filter PHI", get_options_with_counts(df_c, 'PHI', 'tab2'), key='phi2'))
            if s2_phi: df_c = df_c[df_c['PHI'].isin(s2_phi)]
        with c3: 
            ignore_cols = ['ZONE', 'TB Unit', 'PHI', 'Episode ID', 'Patient Name', 'Facility Type', 'Diagnosis Date', 'Initiation Date', 'Outcome Date']
            s2_ind = st.multiselect("Filter by Report Type", [c for c in df_c.columns if c not in ignore_cols], key='ind2')
            s2_stat = st.multiselect("Filter by Status", ["🔴 NEW", "🟢 RESOLVED", "🟡 PERSISTENT"], key='stat2')
        
        cd1, cd2, cd3 = st.columns(3)
        with cd1: diag_dt2 = st.date_input("Diagnosis Date Range", value=[], key="d1_2")
        with cd2: init_dt2 = st.date_input("Initiation Date Range", value=[], key="d2_2")
        with cd3: out_dt2 = st.date_input("Outcome Date Range", value=[], key="d3_2")

        if len(diag_dt2) == 2: df_c = df_c[pd.to_datetime(df_c.get('Diagnosis Date'), errors='coerce').dt.date.between(diag_dt2[0], diag_dt2[1])]
        if len(init_dt2) == 2: df_c = df_c[pd.to_datetime(df_c.get('Initiation Date'), errors='coerce').dt.date.between(init_dt2[0], init_dt2[1])]
        if len(out_dt2) == 2: df_c = df_c[pd.to_datetime(df_c.get('Outcome Date'), errors='coerce').dt.date.between(out_dt2[0], out_dt2[1])]

    if s2_ind or s2_stat:
        mask = pd.Series(False, index=df_c.index)
        for ind in (s2_ind if s2_ind else [c for c in df_c.columns if c not in ignore_cols]):
            if ind in df_c.columns: mask = mask | df_c[ind].isin(s2_stat if s2_stat else ["🔴 NEW", "🟢 RESOLVED", "🟡 PERSISTENT"])
        df_c = df_c[mask]
    
    ind_cols_in_df = [c for c in df_c.columns if c not in ignore_cols]
    new_c = (df_c[ind_cols_in_df] == "🔴 NEW").sum().sum() if ind_cols_in_df else 0
    res_c = (df_c[ind_cols_in_df] == "🟢 RESOLVED").sum().sum() if ind_cols_in_df else 0
    per_c = (df_c[ind_cols_in_df] == "🟡 PERSISTENT").sum().sum() if ind_cols_in_df else 0
    
    st.markdown("##### 📈 Daily Action Status")
    cc1, cc2, cc3, cc4 = st.columns(4)
    with cc1: st.markdown(draw_card("TOTAL PENDENCY", new_c + per_c, "#1f618d", "📄"), unsafe_allow_html=True)
    with cc2: st.markdown(draw_card("🔴 NEW", new_c, "#E74C3C", "🚨"), unsafe_allow_html=True)
    with cc3: st.markdown(draw_card("🟡 PERSISTENT", per_c, "#F1C40F", "⏳"), unsafe_allow_html=True)
    with cc4: st.markdown(draw_card("🟢 RESOLVED", res_c, "#27AE60", "✅"), unsafe_allow_html=True)
    
    # 🎯 FIX: Removing dates from Tab 2 UI display
    df_c_display = df_c.drop(columns=['Diagnosis Date', 'Initiation Date', 'Outcome Date'], errors='ignore')
    st.dataframe(df_c_display, use_container_width=True, hide_index=True)
    if not df_c_display.empty:
        st.download_button("📥 Download Comparison Matrix", convert_df_to_excel(df_c_display, "Comparison_Matrix"), "Comparison.xlsx", key='dl2')

    # ==========================================
    # 📊 DYNAMIC ZONE-WISE PENDENCY TRACKER (ADDED AT THE BOTTOM)
    # ==========================================
    st.write("---")
    st.markdown("<h4 style='color: #1e293b; font-weight: 700;'>📊 Zone-wise Pendency Tracker</h4>", unsafe_allow_html=True)
    
    if not df_c.empty and ind_cols_in_df and 'ZONE' in df_c.columns:
        zone_data = []
        # ERROR FIX: pd.notna(z) ખરાબ ડેટા કે ખાલી ખાના ને લીધે આવતી TypeError ને રોકે છે.
        valid_zones = [z for z in df_c['ZONE'].unique() if pd.notna(z)]
        
        for z in sorted(valid_zones):
            df_z = df_c[df_c['ZONE'] == z]
            n_val = (df_z[ind_cols_in_df] == "🔴 NEW").sum().sum()
            p_val = (df_z[ind_cols_in_df] == "🟡 PERSISTENT").sum().sum()
            r_val = (df_z[ind_cols_in_df] == "🟢 RESOLVED").sum().sum()
            
            zone_data.append({
                'ZONE': z,
                'TOTAL PENDACY (NEW+PERSISTENT)': n_val + p_val,
                'PERSIST': p_val,
                'NEW': n_val,
                'RESOLVE': r_val
            })
        
        df_pendency = pd.DataFrame(zone_data)
        
        if not df_pendency.empty:
            # Excel જેવો કલર સ્કેલ (વધારે પેન્ડન્સી = ઘાટો લાલ કલર)
            styled_df = df_pendency.style.background_gradient(
                subset=['TOTAL PENDACY (NEW+PERSISTENT)'], 
                cmap='Reds' 
            ).format(precision=0)
            
            st.dataframe(styled_df, use_container_width=True, hide_index=True)

# ==========================================
# 🟢 TAB 4: PPT GENERATOR (SMART + CORPORATE + NAAT)
# ==========================================
with tab4:
    st.markdown("<h3 style='text-align: center; color: #27AE60;'>🚀 Enterprise PPT Report Generator</h3>", unsafe_allow_html=True)
    
    with st.container():
        c1, c2, c3 = st.columns(3)
        with c1:
            all_inds = ["Outcome", "UDST", "Not Put On", "SLPA", "Consent", "ADT", "RBS", "ART", "CPT", "HIV"]
            sel_report = st.selectbox("📌 1. Select Report Type", all_inds)
            st.markdown("<div style='background-color:#e8f4f8; padding:10px; border-radius:5px;'><b>📅 Period 1 (Current)</b></div>", unsafe_allow_html=True)
            p1_name = st.text_input("Name for Period 1", "Q1 - 2026")
            p1_diag = st.date_input("Diagnosis Date (P1)", value=[])
            p1_init = st.date_input("Treatment Start Date (P1)", value=[])
            p1_out = st.date_input("Outcome Date (P1)", value=[])
        with c2:
            st.write("")
            st.write("")
            compare_mode = st.checkbox("📊 Enable Comparison (Period 2)")
            if compare_mode:
                st.markdown("<div style='background-color:#fef5e7; padding:10px; border-radius:5px;'><b>📅 Period 2 (Previous)</b></div>", unsafe_allow_html=True)
                p2_name = st.text_input("Name for Period 2", "Q2 - 2026")
                p2_diag = st.date_input("Diagnosis Date (P2)", value=[])
                p2_init = st.date_input("Treatment Start Date (P2)", value=[])
                p2_out = st.date_input("Outcome Date (P2)", value=[])
            else:
                p2_name = "None"
                p2_diag, p2_init, p2_out = [], [], []
        with c3:
            st.markdown("<div style='background-color:#e9ecef; padding:10px; border-radius:5px;'><b>🎨 3. Presentation Rules</b></div>", unsafe_allow_html=True)
            st.write("")
            color_rule = st.radio("Color Scale Rules:", ["High is Bad (Red) 🔴", "High is Good (Green) 🟢"])
            high_is_bad = True if "Bad" in color_rule else False
            if compare_mode: color_target = st.radio("Apply Color Formatting On:", [p1_name, p2_name, "Grand Total"])
            else: color_target = p1_name

    def apply_date_filters(df, diag, init, out):
        mask = pd.Series(True, index=df.index)
        if len(diag) == 2: mask &= pd.to_datetime(df.get('Diagnosis Date'), errors='coerce').dt.date.between(diag[0], diag[1])
        if len(init) == 2: mask &= pd.to_datetime(df.get('Initiation Date'), errors='coerce').dt.date.between(init[0], init[1])
        if len(out) == 2: mask &= pd.to_datetime(df.get('Outcome Date'), errors='coerce').dt.date.between(out[0], out[1])
        return mask

    def generate_smart_ppt(df, report_name):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError: return None, "⚠️ PPTX લાઈબ્રેરી ઇન્સ્ટોલ નથી!"

        prs = Presentation()
        m1 = apply_date_filters(df, p1_diag, p1_init, p1_out)
        m1 &= df.get('Pending Status', pd.Series(dtype=str)).astype(str).str.contains(report_name, na=False)
        df_p1 = df[m1].copy()

        df_p2 = pd.DataFrame()
        if compare_mode:
            m2 = apply_date_filters(df, p2_diag, p2_init, p2_out)
            m2 &= df.get('Pending Status', pd.Series(dtype=str)).astype(str).str.contains(report_name, na=False)
            df_p2 = df[m2].copy()

        def get_bg_color(val, max_val):
            if max_val == 0 or pd.isna(val) or val == 0: return RGBColor(255, 255, 255)
            ratio = val / max_val
            if not high_is_bad: ratio = 1 - ratio 
            if ratio > 0.66: return RGBColor(241, 148, 138)
            elif ratio > 0.33: return RGBColor(249, 231, 159)
            else: return RGBColor(171, 235, 198)

        def add_slide_table(title_text, curr_df, prev_df, entity_col_name):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if os.path.exists("images/amc.png"): slide.shapes.add_picture("images/amc.png", Inches(0.2), Inches(0.15), width=Inches(0.7))
            if os.path.exists("images/ntep.jpg"): slide.shapes.add_picture("images/ntep.jpg", Inches(9.1), Inches(0.15), width=Inches(0.7))
            title = slide.shapes.title
            title.text = title_text
            title.top = Inches(0.25); title.left = Inches(1.2); title.width = Inches(7.6)
            title.text_frame.paragraphs[0].font.size = Pt(26)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
            
            if curr_df.empty and prev_df.empty: return
            if compare_mode:
                final_df = pd.merge(curr_df, prev_df, on=entity_col_name, how='outer').fillna(0)
                final_df['Grand Total'] = final_df[p1_name] + final_df[p2_name]
                final_df = final_df.sort_values(by='Grand Total', ascending=False)
                col_names = [entity_col_name, p1_name, p2_name, 'Grand Total']
            else:
                final_df = curr_df.sort_values(by=p1_name, ascending=False)
                col_names = [entity_col_name, p1_name]

            rows = len(final_df) + 1
            cols = len(col_names)
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.3), Inches(8.4), Inches(0.4))
            table = table_shape.table
            if cols == 2: table.columns[0].width = Inches(5.4); table.columns[1].width = Inches(3.0)
            elif cols == 4: table.columns[0].width = Inches(4.0); table.columns[1].width = Inches(1.5); table.columns[2].width = Inches(1.5); table.columns[3].width = Inches(1.4)
            for i, c_name in enumerate(col_names):
                cell = table.cell(0, i)
                cell.text = c_name
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
                if i > 0: cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            target_idx = col_names.index(color_target)
            max_value = final_df.iloc[:, target_idx].max() if not final_df.empty else 0
            for i, (_, row) in enumerate(final_df.iterrows()):
                name_val = str(row[entity_col_name])
                table.cell(i+1, 0).text = name_val
                c1_p = table.cell(i+1, 1).text_frame.paragraphs[0]
                c1_p.text = str(int(row[p1_name]))
                c1_p.alignment = PP_ALIGN.CENTER
                if cols == 4:
                    c2_p = table.cell(i+1, 2).text_frame.paragraphs[0]
                    c2_p.text = str(int(row[p2_name]))
                    c2_p.alignment = PP_ALIGN.CENTER
                    c3_p = table.cell(i+1, 3).text_frame.paragraphs[0]
                    c3_p.text = str(int(row['Grand Total']))
                    c3_p.alignment = PP_ALIGN.CENTER
                if "PRIVATE FACILITIES" in name_val:
                    for j in range(cols):
                        c = table.cell(i+1, j)
                        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(235, 237, 239)
                        c.text_frame.paragraphs[0].font.bold = True
                else:
                    c_target = table.cell(i+1, target_idx)
                    val_target = row.iloc[target_idx]
                    c_target.fill.solid()
                    c_target.fill.fore_color.rgb = get_bg_color(val_target, max_value)

        def get_summary(temp_df, group_col, val_name):
            if temp_df.empty: return pd.DataFrame(columns=[group_col, val_name])
            if group_col == 'PHI':
                pub_mask = temp_df.get('Facility Type', pd.Series(dtype=str)).astype(str).str.upper().isin(['PUBLIC', 'PHI'])
                pub_sum = temp_df[pub_mask].groupby('PHI').size().reset_index(name=val_name)
                priv_count = len(temp_df[~pub_mask])
                if priv_count > 0:
                    priv_row = pd.DataFrame({'PHI': ['PRIVATE FACILITIES (TOTAL)'], val_name: [priv_count]})
                    return pd.concat([pub_sum, priv_row], ignore_index=True)
                return pub_sum
            else: return temp_df.groupby(group_col).size().reset_index(name=val_name)

        if st.session_state.role == "ZONE":
            tu_curr = get_summary(df_p1, 'TB Unit', p1_name)
            tu_prev = get_summary(df_p2, 'TB Unit', p2_name) if compare_mode else pd.DataFrame()
            add_slide_table(f"{st.session_state.target} Zone - {sel_report} Pending", tu_curr, tu_prev, 'TB Unit')
            tus = sorted(pd.concat([df_p1.get('TB Unit', pd.Series()), df_p2.get('TB Unit', pd.Series()) if compare_mode else pd.Series()]).dropna().unique())
            for tu in tus:
                phi_curr = get_summary(df_p1[df_p1.get('TB Unit') == tu], 'PHI', p1_name)
                phi_prev = get_summary(df_p2[df_p2.get('TB Unit') == tu], 'PHI', p2_name) if compare_mode else pd.DataFrame()
                add_slide_table(f"TU: {tu} - {sel_report} Pending", phi_curr, phi_prev, 'PHI')

        elif st.session_state.role == "ADMIN":
            z_curr = get_summary(df_p1, 'ZONE', p1_name)
            z_prev = get_summary(df_p2, 'ZONE', p2_name) if compare_mode else pd.DataFrame()
            add_slide_table(f"All Zones - {sel_report} Pending", z_curr, z_prev, 'ZONE')
            zones = sorted(pd.concat([df_p1.get('ZONE', pd.Series()), df_p2.get('ZONE', pd.Series()) if compare_mode else pd.Series()]).dropna().unique())
            for z in zones:
                phi_curr = get_summary(df_p1[df_p1.get('ZONE') == z], 'PHI', p1_name)
                phi_prev = get_summary(df_p2[df_p2.get('ZONE') == z], 'PHI', p2_name) if compare_mode else pd.DataFrame()
                add_slide_table(f"Zone: {z} - {sel_report} Pending", phi_curr, phi_prev, 'PHI')
        else:
            phi_curr = get_summary(df_p1, 'PHI', p1_name)
            phi_prev = get_summary(df_p2, 'PHI', p2_name) if compare_mode else pd.DataFrame()
            add_slide_table(f"{st.session_state.target} - {sel_report} Pending", phi_curr, df_p2, 'PHI')

        out_io = io.BytesIO()
        prs.save(out_io)
        return out_io.getvalue(), "Success"

    st.markdown("<br>", unsafe_allow_html=True)
    if st.button("✨ Generate Custom PPT ✨", use_container_width=True):
        with st.spinner("Generating beautiful Enterprise PPT slides... Please wait..."):
            ppt_bytes, status = generate_smart_ppt(df_master, sel_report)
            if ppt_bytes:
                st.success("✅ PPT 100% તૈયાર છે! નીચેના બટન પર ક્લિક કરીને ડાઉનલોડ કરો.")
                st.download_button(label=f"📥 Download {sel_report}_Analysis.pptx", data=ppt_bytes, file_name=f"{sel_report}_Analysis.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            else: st.error(status)


    # ==========================================
    # 🎯 2. MNC CORPORATE TARGET ACHIEVEMENT DECK
    # ==========================================
    st.markdown("<br><hr style='margin: 30px 0; border: 2px solid #e8f4f8;'>", unsafe_allow_html=True)
    st.markdown("<h3 style='text-align: center; color: #2C3E50;'>📈 Corporate Performance Deck (Zone + UHC/CHC/HOSPITAL)</h3>", unsafe_allow_html=True)

    with st.container():
        tc1, tc2, tc3 = st.columns(3)
        with tc1:
            st.markdown("<div style='background-color:#fef9e7; padding:10px; border-radius:5px;'><b>🗓️ 1. Date Selection</b></div>", unsafe_allow_html=True)
            target_dates = st.date_input("Select Dates to Sum (e.g., June 1 to June 5)", value=[], key="t_dates")
        with tc2:
            st.markdown("<div style='background-color:#e8f8f5; padding:10px; border-radius:5px;'><b>🔢 2. Target Multiplier</b></div>", unsafe_allow_html=True)
            working_days = st.number_input("Enter Total Working Days", min_value=1, max_value=31, value=5, key="t_wdays")
        with tc3:
            st.markdown("<div style='background-color:#ebedf0; padding:10px; border-radius:5px;'><b>⚙️ 3. Action</b></div>", unsafe_allow_html=True)
            st.write("")
            btn_generate_target = st.button("✨ Generate Full Deck ✨", use_container_width=True)

    def generate_corporate_target_ppt(selected_dates, w_days):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            import re
            
            def add_corporate_slide(prs_obj, title_text):
                slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
                if os.path.exists("images/amc.png"): slide.shapes.add_picture("images/amc.png", Inches(0.2), Inches(0.15), width=Inches(0.6))
                if os.path.exists("images/ntep.jpg"): slide.shapes.add_picture("images/ntep.jpg", Inches(9.2), Inches(0.15), width=Inches(0.6))
                title = slide.shapes.title; title.text = title_text
                title.top = Inches(0.25); title.left = Inches(1.0)
                title.width = Inches(8.0); title.height = Inches(0.8)
                title.text_frame.paragraphs[0].font.size = Pt(24); title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
                return slide

            def format_corporate_table(table_obj, df_data, col_widths, font_size=12):
                rows, cols = len(df_data) + 1, len(df_data.columns)
                for i, width in enumerate(col_widths): table_obj.columns[i].width = width
                for i, col_name in enumerate(df_data.columns):
                    cell = table_obj.cell(0, i); cell.text = col_name
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.size = Pt(12)
                    if i > 1: cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                for i in range(1, rows):
                    for j in range(cols):
                        cell = table_obj.cell(i, j)
                        for p in cell.text_frame.paragraphs: 
                            p.font.size = Pt(font_size)
                            if j > 1: p.alignment = PP_ALIGN.CENTER

            def extract_num(val):
                nums = re.findall(r'^(\d+)', str(val).strip())
                return int(nums[0]) if nums else 0

            def get_multi_color(pct):
                if pct >= 100: return RGBColor(46, 204, 113)
                elif pct >= 75: return RGBColor(171, 235, 198)
                elif pct >= 50: return RGBColor(249, 231, 159)
                elif pct >= 25: return RGBColor(245, 176, 65)
                else: return RGBColor(231, 76, 60)

            if len(selected_dates) == 2:
                date_list = pd.date_range(start=selected_dates[0], end=selected_dates[1]).tolist()
                target_date_strings = [f"{d.strftime('%b')} {d.day}, {d.year}" for d in date_list]
            else: return None, "⚠️ Please select a start and end date."

            prs = Presentation()
            fixed_targets = {"Central": 59, "North": 122, "East": 117, "South": 159, "West": 121, "North West": 77, "South West": 55, "AMC": 710}
            
            zone_urls = [
                "https://docs.google.com/spreadsheets/d/19Whbn-0bGNxVcxiGmp9fCq44dKeNZXAAbPiXtVf3zcs/export?format=csv&gid=972568835", # JUNE ZONE
                "https://docs.google.com/spreadsheets/d/19Whbn-0bGNxVcxiGmp9fCq44dKeNZXAAbPiXtVf3zcs/export?format=csv&gid=470337901"  # MAY ZONE
            ]
            
            fac_urls = [
                "https://docs.google.com/spreadsheets/d/19Whbn-0bGNxVcxiGmp9fCq44dKeNZXAAbPiXtVf3zcs/export?format=csv&gid=0", # JUNE FACILITY
                "https://docs.google.com/spreadsheets/d/19Whbn-0bGNxVcxiGmp9fCq44dKeNZXAAbPiXtVf3zcs/export?format=csv&gid=218126721" # MAY FACILITY
            ]

            # ----------------------------------------------------
            # 1️⃣ AGGREGATE ZONE DATA ACROSS ALL SHEETS
            # ----------------------------------------------------
            zone_achievements = {z: 0 for z in fixed_targets.keys() if z != "AMC"}
            
            for url in zone_urls:
                try:
                    df_sheet1 = pd.read_csv(url, header=None)
                    h_idx1 = 0
                    for i in range(3):
                        if any(td in df_sheet1.iloc[i].fillna("").astype(str).tolist() for td in target_date_strings):
                            h_idx1 = i; break
                            
                    header_row1 = df_sheet1.iloc[h_idx1].fillna("").astype(str)
                    col_indices1 = [idx for idx, val in enumerate(header_row1) if val.replace("  ", " ").strip() in target_date_strings]
                    
                    if col_indices1:
                        for row_idx in range(h_idx1 + 1, len(df_sheet1)):
                            z_name = str(df_sheet1.iloc[row_idx, 0]).strip().title()
                            if z_name in zone_achievements:
                                ach_total = sum([extract_num(df_sheet1.iloc[row_idx, c]) for c in col_indices1])
                                zone_achievements[z_name] += ach_total
                except: continue

            # Build Zone Table
            res1 = []
            for z_name, ach_total in zone_achievements.items():
                if ach_total > 0 or True:
                    t_day = fixed_targets[z_name]
                    m_target = t_day * w_days
                    pct = round((ach_total / m_target) * 100, 1) if m_target > 0 else 0
                    res1.append({"ZONE": z_name, "TARGET PER DAY": t_day, "MONTH TARGET": m_target, "TOTAL ACHIEVED": ach_total, "ACHIEVEMENT %": pct})
            
            if res1:
                amc_target_day = fixed_targets["AMC"]
                amc_month_target = amc_target_day * w_days
                amc_achieved = sum(r["TOTAL ACHIEVED"] for r in res1)
                amc_pct = round((amc_achieved / amc_month_target) * 100, 1) if amc_month_target > 0 else 0
                res1.append({"ZONE": "AMC", "TARGET PER DAY": amc_target_day, "MONTH TARGET": amc_month_target, "TOTAL ACHIEVED": amc_achieved, "ACHIEVEMENT %": amc_pct})

                df_display1 = pd.DataFrame(res1)
                df_display1["ACHIEVEMENT %"] = df_display1["ACHIEVEMENT %"].astype(str) + "%"
                
                s1 = add_corporate_slide(prs, f"Cumulative Target Achievement (Days: {w_days})")
                t1_shape = s1.shapes.add_table(len(df_display1) + 1, len(df_display1.columns), Inches(0.5), Inches(1.3), Inches(9.0), Inches(0.4))
                format_corporate_table(t1_shape.table, df_display1, [Inches(2.0), Inches(1.5), Inches(1.8), Inches(1.7), Inches(2.0)])
                
                for i, row in df_display1.iterrows():
                    for j in range(len(df_display1.columns)):
                        cell = t1_shape.table.cell(i+1, j)
                        cell.text = str(row.iloc[j])
                        if row['ZONE'] == "AMC":
                            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(235, 237, 239)
                            cell.text_frame.paragraphs[0].font.bold = True
                        if j == 4 and row['ZONE'].upper() != "AMC":
                            pct_val = float(str(row.iloc[j]).replace('%', ''))
                            cell.fill.solid(); cell.fill.fore_color.rgb = get_multi_color(pct_val)

            # ----------------------------------------------------
            # 2️⃣ AGGREGATE FACILITY DATA ACROSS ALL SHEETS (Including Hospitals)
            # ----------------------------------------------------
            fac_achievements = {}

            for url in fac_urls:
                try:
                    df_fac = pd.read_csv(url, header=None)
                    h_idx2 = 0
                    for i in range(4):
                        if any(td in df_fac.iloc[i].fillna("").astype(str).tolist() for td in target_date_strings):
                            h_idx2 = i; break
                    
                    header_fac = df_fac.iloc[h_idx2].fillna("").astype(str)
                    col_indices_fac = [idx for idx, val in enumerate(header_fac) if val.replace("  ", " ").strip() in target_date_strings]

                    if col_indices_fac:
                        for row_idx in range(h_idx2 + 1, len(df_fac)):
                            zone_guj = str(df_fac.iloc[row_idx, 0]).strip()
                            fac_name = str(df_fac.iloc[row_idx, 1]).strip()
                            if "કુલ" in fac_name or "કુલ" in zone_guj or fac_name in ["", "nan", "None"]: continue
                                
                            achieved_total = sum([extract_num(df_fac.iloc[row_idx, c]) for c in col_indices_fac])
                            fac_type = "OTHER"
                            
                            if "અર્બન હેલ્થ સેન્ટર" in fac_name: fac_type = "UHC"
                            elif "સામુહીક" in fac_name or "સામુહિક" in fac_name: fac_type = "CHC"
                            elif "હોસ્પિટલ" in fac_name: fac_type = "HOSPITAL"
                            
                            if fac_type in ["UHC", "CHC", "HOSPITAL"]:
                                dict_key = (zone_guj, fac_name, fac_type)
                                if dict_key in fac_achievements: fac_achievements[dict_key] += achieved_total
                                else: fac_achievements[dict_key] = achieved_total
                except: continue

            # Build Facility Tables
            if fac_achievements:
                fac_data = []
                for (zone_guj, fac_name, fac_type), achieved_total in fac_achievements.items():
                    
                    if fac_type == "UHC": target_daily = 4
                    elif fac_type == "CHC": target_daily = 16
                    elif fac_type == "HOSPITAL": target_daily = 30 
                    
                    month_target = target_daily * w_days
                    ach_pct = round((achieved_total / month_target) * 100, 1) if month_target > 0 else 0
                    fac_data.append({"Zone": zone_guj, "Facility Name": fac_name, "Type": fac_type, "Target": month_target, "Achieved": achieved_total, "Achievement %": ach_pct})
                
                df_fac_processed = pd.DataFrame(fac_data)

                # --- 📉 UHC SLIDES ---
                if not df_fac_processed.empty:
                    df_uhc = df_fac_processed[(df_fac_processed["Type"] == "UHC") & (df_fac_processed["Achievement %"] < 75)].sort_values("Achievement %").drop(columns=["Type"]).reset_index(drop=True)
                    df_uhc_display = df_uhc.copy()
                    df_uhc_display["Achievement %"] = df_uhc_display["Achievement %"].astype(str) + "%"
                    
                    chunk_size = 12
                    for i in range(0, len(df_uhc_display), chunk_size):
                        chunk = df_uhc_display.iloc[i:i+chunk_size]
                        s2 = add_corporate_slide(prs, f"📉 UHCs Requiring Attention (< 75%){' (Part ' + str(i//chunk_size + 1) + ')' if len(df_uhc_display)>chunk_size else ''}")
                        t2 = s2.shapes.add_table(len(chunk) + 1, len(chunk.columns), Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.35))
                        format_corporate_table(t2.table, chunk, [Inches(1.5), Inches(4.0), Inches(1.0), Inches(1.0), Inches(1.5)], font_size=11)
                        for row_idx_c, (orig_idx, row) in enumerate(chunk.iterrows()):
                            for j in range(len(chunk.columns)):
                                cell = t2.table.cell(row_idx_c+1, j); cell.text = str(row.iloc[j])
                                for p in cell.text_frame.paragraphs: 
                                    p.font.size = Pt(11)
                                    if j > 1: p.alignment = PP_ALIGN.CENTER
                                if j == 4:
                                    cell.fill.solid(); cell.fill.fore_color.rgb = get_multi_color(df_uhc.iloc[orig_idx]["Achievement %"])

                # --- 🏥 CHC SLIDES ---
                if not df_fac_processed.empty:
                    df_chc = df_fac_processed[df_fac_processed["Type"] == "CHC"].sort_values("Achievement %", ascending=False).drop(columns=["Type"]).reset_index(drop=True)
                    df_chc_display = df_chc.copy()
                    df_chc_display["Achievement %"] = df_chc_display["Achievement %"].astype(str) + "%"
                    
                    for i in range(0, len(df_chc_display), 12):
                        chunk = df_chc_display.iloc[i:i+12]
                        s3 = add_corporate_slide(prs, f"🏥 All CHCs Performance Overview{' (Part ' + str(i//12 + 1) + ')' if len(df_chc_display)>12 else ''}")
                        t3 = s3.shapes.add_table(len(chunk) + 1, len(chunk.columns), Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.35))
                        format_corporate_table(t3.table, chunk, [Inches(1.5), Inches(4.0), Inches(1.0), Inches(1.0), Inches(1.5)], font_size=11)
                        for row_idx_c, (orig_idx, row) in enumerate(chunk.iterrows()):
                            for j in range(len(chunk.columns)):
                                cell = t3.table.cell(row_idx_c+1, j); cell.text = str(row.iloc[j])
                                for p in cell.text_frame.paragraphs: 
                                    p.font.size = Pt(11)
                                    if j > 1: p.alignment = PP_ALIGN.CENTER
                                if j == 4:
                                    cell.fill.solid(); cell.fill.fore_color.rgb = get_multi_color(df_chc.iloc[orig_idx]["Achievement %"])
                                elif row_idx_c % 2 != 0: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(242, 243, 244)
                
                # --- 🏥 HOSPITAL SLIDES ---
                if not df_fac_processed.empty:
                    df_hosp = df_fac_processed[df_fac_processed["Type"] == "HOSPITAL"].sort_values("Achievement %", ascending=False).drop(columns=["Type"]).reset_index(drop=True)
                    
                    if not df_hosp.empty:
                        df_hosp_display = df_hosp.copy()
                        df_hosp_display["Achievement %"] = df_hosp_display["Achievement %"].astype(str) + "%"
                        
                        for i in range(0, len(df_hosp_display), 12):
                            chunk = df_hosp_display.iloc[i:i+12]
                            s4 = add_corporate_slide(prs, f"🏥 Hospital Performance Overview{' (Part ' + str(i//12 + 1) + ')' if len(df_hosp_display)>12 else ''}")
                            t4 = s4.shapes.add_table(len(chunk) + 1, len(chunk.columns), Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.35))
                            format_corporate_table(t4.table, chunk, [Inches(1.5), Inches(4.0), Inches(1.0), Inches(1.0), Inches(1.5)], font_size=11)
                            for row_idx_c, (orig_idx, row) in enumerate(chunk.iterrows()):
                                for j in range(len(chunk.columns)):
                                    cell = t4.table.cell(row_idx_c+1, j); cell.text = str(row.iloc[j])
                                    for p in cell.text_frame.paragraphs: 
                                        p.font.size = Pt(11)
                                        if j > 1: p.alignment = PP_ALIGN.CENTER
                                    if j == 4:
                                        cell.fill.solid(); cell.fill.fore_color.rgb = get_multi_color(df_hosp.iloc[orig_idx]["Achievement %"])
                                    elif row_idx_c % 2 != 0: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(242, 243, 244)
            
            out_io = io.BytesIO()
            prs.save(out_io)
            return out_io.getvalue(), "Success"
        except Exception as e: return None, f"⚠️ Error: {str(e)}"

    if btn_generate_target:
        if len(target_dates) != 2: st.error("⚠️ Please select both a Start Date and End Date.")
        else:
            with st.spinner("Fetching Live Sheet Data and generating Corporate Deck..."):
                target_ppt_bytes, t_status = generate_corporate_target_ppt(target_dates, working_days)
                if target_ppt_bytes:
                    st.success("✅ Corporate Presentation Deck Ready!")
                    st.download_button(label="📥 Download Corporate_Deck.pptx", data=target_ppt_bytes, file_name="Corporate_Performance_Deck.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key="dl_target_ppt")
                else: st.error(t_status)


    # ==========================================
    # 🎯 3. NAAT UTILIZATION REPORT DECK (WITH 100% BULLETPROOF FIXES)
    # ==========================================
    st.markdown("<br><hr style='margin: 30px 0; border: 2px solid #e8f4f8;'>", unsafe_allow_html=True)
    st.markdown("<h3 style='text-align: center; color: #E67E22;'>🔬 NAAT Utilization Report Generator</h3>", unsafe_allow_html=True)

    with st.container():
        nc1, nc2, nc3 = st.columns(3)
        with nc1:
            st.markdown("<div style='background-color:#fef9e7; padding:10px; border-radius:5px;'><b>🗓️ 1. Date Selection</b></div>", unsafe_allow_html=True)
            naat_dates = st.date_input("Select NAAT Dates", value=[], key="n_dates")
        with nc2:
            st.markdown("<div style='background-color:#e8f8f5; padding:10px; border-radius:5px;'><b>🔢 2. Divisor</b></div>", unsafe_allow_html=True)
            naat_wdays = st.number_input("Enter Working Days (for Average)", min_value=1, max_value=31, value=5, key="n_wdays")
        with nc3:
            st.markdown("<div style='background-color:#ebedf0; padding:10px; border-radius:5px;'><b>⚙️ 3. Action</b></div>", unsafe_allow_html=True)
            st.write("")
            btn_generate_naat = st.button("✨ Generate NAAT PPT ✨", use_container_width=True)

    def generate_naat_ppt(selected_dates, w_days):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            def add_corporate_slide(prs_obj, title_text):
                slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
                if os.path.exists("images/amc.png"): slide.shapes.add_picture("images/amc.png", Inches(0.2), Inches(0.15), width=Inches(0.6))
                if os.path.exists("images/ntep.jpg"): slide.shapes.add_picture("images/ntep.jpg", Inches(9.2), Inches(0.15), width=Inches(0.6))
                title = slide.shapes.title; title.text = title_text
                title.top = Inches(0.25); title.left = Inches(1.0)
                title.width = Inches(8.0); title.height = Inches(0.8)
                title.text_frame.paragraphs[0].font.size = Pt(24); title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
                return slide

            def format_corporate_table(table_obj, df_data, col_widths, font_size=12):
                rows, cols = len(df_data) + 1, len(df_data.columns)
                for i, width in enumerate(col_widths): table_obj.columns[i].width = width
                for i, col_name in enumerate(df_data.columns):
                    cell = table_obj.cell(0, i); cell.text = col_name
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.size = Pt(12)
                    if i > 1: cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                for i in range(1, rows):
                    for j in range(cols):
                        cell = table_obj.cell(i, j)
                        for p in cell.text_frame.paragraphs: 
                            p.font.size = Pt(font_size)
                            if j > 1: p.alignment = PP_ALIGN.CENTER
            
            if len(selected_dates) == 2:
                date_list = pd.date_range(start=selected_dates[0], end=selected_dates[1]).tolist()
            else: return None, "⚠️ Please select a start and end date."

            # 🎯 EXPLICIT MONTH MAPPING
            naat_urls = {
                6: "https://docs.google.com/spreadsheets/d/1a1F3BZsGjgM8-_JY0ohbvsODxM6cPPLksDRFlaVgB0s/export?format=csv&gid=718682714", # JUNE
                5: "https://docs.google.com/spreadsheets/d/1a1F3BZsGjgM8-_JY0ohbvsODxM6cPPLksDRFlaVgB0s/export?format=csv&gid=910963940"  # MAY
            }
            
            site_totals = {}
            found_any_date = False

            # Group the selected dates by their actual month
            dates_by_month = {}
            for d in date_list:
                if d.month not in dates_by_month: dates_by_month[d.month] = []
                dates_by_month[d.month].append(d)
                
            for target_month, m_dates in dates_by_month.items():
                if target_month not in naat_urls: continue 
                    
                try:
                    df_naat = pd.read_csv(naat_urls[target_month], header=None)
                    
                    # 🛡️ THE FIX: Forward-fill Column 0 so UCHC THALTEJ persists down its merged cells!
                    df_naat[0] = df_naat[0].replace(["", "nan", "NaN", "None"], pd.NA).ffill()
                    
                    # 🛡️ DYNAMIC ROW FINDER
                    date_row_idx = 0
                    header_row_idx = 1
                    for i in range(min(5, len(df_naat))):
                        row_str = " ".join(df_naat.iloc[i].fillna("").astype(str).values).lower()
                        if "tested" in row_str or "sample sent" in row_str:
                            header_row_idx = i
                            date_row_idx = i - 1
                            break
                            
                    # 🛡️ Forward-fill the Date Row so sub-columns inherit the correct date
                    date_series = df_naat.iloc[date_row_idx].replace(["", "nan", "NaN", "None"], pd.NA).ffill()
                    header_series = df_naat.iloc[header_row_idx].fillna("").astype(str).str.upper().str.strip()
                    
                    parsed_dates = pd.to_datetime(date_series, errors='coerce').dt.date
                    date_str_series = date_series.astype(str).str.lower().str.replace("-", "").str.replace("/", "").str.replace(" ", "")

                    # Slice to data rows
                    df_valid = df_naat.iloc[header_row_idx + 1:].copy()
                    
                    # 🛡️ Ignore the TOTAL rows dynamically to prevent double counting
                    mask_tot = (df_valid[0].astype(str).str.upper().str.contains("TOTAL", na=False) | 
                                df_valid[1].astype(str).str.upper().str.contains("TOTAL", na=False) | 
                                df_valid[2].astype(str).str.upper().str.contains("TOTAL", na=False))
                    df_valid = df_valid[~mask_tot]
                    
                    for d in m_dates:
                        d_date = d.date()
                        f1 = f"{d.month:02d}{d.day:02d}" 
                        f2 = f"{d.day:02d}{d.month:02d}" 
                        f3 = f"{d.month}{d.day}"         
                        f4 = d.strftime("%d%b").lower()  
                        f5 = d.strftime("%b%d").lower()  
                        
                        sheet_tested_cols = []
                        
                        for i in range(len(date_series)):
                            is_match = False
                            if pd.notna(parsed_dates[i]) and parsed_dates[i] == d_date:
                                is_match = True
                            else:
                                v_str = date_str_series[i]
                                if f1 in v_str or f2 in v_str or f3 in v_str or f4 in v_str or f5 in v_str:
                                    is_match = True
                                    
                            # 🛡️ We strictly only pull the TESTED column!
                            if is_match and "TEST" in header_series[i]:
                                sheet_tested_cols.append(i)
                                break 
                        
                        if sheet_tested_cols:
                            found_any_date = True
                            col_idx = sheet_tested_cols[0]
                            
                            for _, row in df_valid.iterrows():
                                site_name = str(row[0]).strip()
                                if site_name not in ["", "nan", "NaN", "None"]:
                                    if col_idx < len(row):
                                        raw_val = row[col_idx]
                                        if pd.notna(raw_val) and str(raw_val).strip() != "":
                                            try: 
                                                sum_val = float(raw_val)
                                                if site_name in site_totals: site_totals[site_name] += sum_val
                                                else: site_totals[site_name] = sum_val
                                            except: pass
                except: continue

            if not found_any_date: return None, "⚠️ Could not process data for the selected dates. Please check the date format in Google Sheets."
            
            grouped = pd.DataFrame(list(site_totals.items()), columns=['NAAT Site', 'Tested'])
            
            def format_avg(val): return int(val) if float(val).is_integer() else round(float(val), 1)
            grouped['Tested'] = grouped['Tested'].astype(int)
            grouped['Average'] = (grouped['Tested'] / w_days).apply(format_avg)
            
            def clean_site(s):
                import re
                c = str(s).upper().replace("CBNAAT", "").replace("TRUNAAT", "").strip(" -,")
                return c if c not in ["NAN", "NONE", ""] else ""
            grouped['NAAT Site'] = grouped['NAAT Site'].apply(clean_site)
            grouped = grouped[grouped['NAAT Site'] != ""]
            
            zone_map_strict = {"MC- CIVIL HOSPITAL, AMC": "Central", "MC-GCS MEDICAL COLLEGE, AMC": "North", "MC GMERS SOLA": "North West", "DH SCL GEN. HOSP.": "North", "UCHC VATVA": "South", "UCHC SABARMATI": "West", "MC-NHL MEDICAL COLLEGE, AMC": "West", "UCHC THALTEJ": "North West", "NARENDRA MODI MC": "South", "FAISALNAGAR CHC": "South", "UCHC DANILIMDA": "South", "UCHC BEHERAMPURA": "South", "CHC VASTRAL": "East", "SDH ESIC MODEL HOSP.": "North", "UHC RANIP": "West", "MC-NARENDRA MODI MEDICAL COLLEGE": "South", "UCHC CHANDKHEDA": "West", "UCHC RAKHIAL": "North", "CHC SARKHEJ": "South West", "UCHC NARODA": "North", "UHC SAIJPUR": "North", "MC-DR. M K SHAH MEDICAL COLLEGE AND RESEARCH CENTER AMC": "West", "UHC SHAHPUR": "Central", "UHC STADIUM": "West", "UHC JAMALPUR": "Central", "UHC GHATLODIA": "North West", "UHC VIRATNAGAR": "East", "UCHC GOMTIPUR": "East", "UHC ISANPUR": "South", "UHC BHAIPURA": "East", "JODHPUR UHC": "South West", "UHC NAVRANGPURA": "West"}
            import re
            clean_zone_map = {re.sub(r'[^A-Z0-9]', '', k.replace("CBNAAT","").replace("TRUNAAT","").upper()): v for k,v in zone_map_strict.items()}

            def get_zone(site):
                c_site = re.sub(r'[^A-Z0-9]', '', str(site).replace("CBNAAT","").replace("TRUNAAT","").upper())
                for k, v in clean_zone_map.items():
                    if k in c_site or c_site in k: return v
                if "SOLA" in c_site: return "North West"
                if "NHL" in c_site or "SABARMATI" in c_site: return "West"
                if "GCS" in c_site or "SHARDABEN" in c_site: return "East"
                if "ASARWA" in c_site or "CIVIL" in c_site: return "Central"
                if "VATVA" in c_site or "MANINAGAR" in c_site or "KANKARIA" in c_site: return "South"
                return "AMC" 
                
            grouped.insert(0, 'Zone', grouped['NAAT Site'].apply(get_zone))
            grouped = grouped.sort_values(by=['Zone', 'Tested'], ascending=[True, False]).reset_index(drop=True)
            
            total_tested = int(grouped['Tested'].sum())
            total_avg = format_avg(total_tested / w_days)
            total_row = pd.DataFrame([{"Zone": "AMC", "NAAT Site": "TOTAL", "Tested": total_tested, "Average": total_avg}])
            grouped = pd.concat([grouped, total_row], ignore_index=True)
            
            prs = Presentation()
            chunk_size = 13
            for i in range(0, len(grouped), chunk_size):
                chunk = grouped.iloc[i:i+chunk_size]
                title_suffix = f" (Part {i//chunk_size + 1})" if len(grouped) > chunk_size else ""
                s = add_corporate_slide(prs, f"🔬 NAAT Utilization Report{title_suffix}")
                t_shape = s.shapes.add_table(len(chunk) + 1, len(chunk.columns), Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.35))
                format_corporate_table(t_shape.table, chunk, [Inches(1.5), Inches(4.5), Inches(1.2), Inches(1.2)], font_size=11)
                
                for row_idx, (orig_idx, row) in enumerate(chunk.iterrows()):
                    is_total_row = (row['Zone'] == "AMC" and row['NAAT Site'] == "TOTAL")
                    for j in range(len(chunk.columns)):
                        cell = t_shape.table.cell(row_idx+1, j); cell.text = str(row.iloc[j])
                        if is_total_row:
                            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(235, 237, 239)
                            for p in cell.text_frame.paragraphs: p.font.bold = True; p.font.size = Pt(11); p.alignment = PP_ALIGN.CENTER if j > 1 else None
                        else:
                            for p in cell.text_frame.paragraphs: p.font.size = Pt(11); p.alignment = PP_ALIGN.CENTER if j > 1 else None
                            if j == 3 and float(row['Average']) < 16: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(241, 148, 138) 
                            elif row_idx % 2 != 0: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(242, 243, 244) 
                            
            out_io = io.BytesIO()
            prs.save(out_io)
            return out_io.getvalue(), "Success"

        except Exception as e: return None, f"⚠️ Error: {str(e)}"

    if btn_generate_naat:
        if len(naat_dates) != 2: st.error("⚠️ Please select both a Start Date and End Date.")
        else:
            with st.spinner("Analyzing NAAT Data and generating PPT..."):
                naat_ppt_bytes, n_status = generate_naat_ppt(naat_dates, naat_wdays)
                if naat_ppt_bytes:
                    st.success("✅ NAAT Utilization Deck Ready!")
                    st.download_button(label="📥 Download NAAT_Report.pptx", data=naat_ppt_bytes, file_name="NAAT_Utilization_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key="dl_naat_ppt")
                else: st.error(n_status)

# ==========================================
# 🟢 TAB 5: DIFFERENTIATED CARE (MINI BOXES, DYNAMIC MATRIX & COMPARISON ENGINE)
# ==========================================
with tab5:
    st.markdown("<h3 style='color: #1f618d;'>🏥 Differentiated Care Tracking System</h3>", unsafe_allow_html=True)
    
    if df_dc_new.empty:
        st.warning("⚠️ ડેટા મળ્યો નથી. ગુગલ શીટ અને લોગિન ઝોન ચેક કરો.")
    else:
        with st.expander("🔽 Filters & Dates (Applies to Current Status)", expanded=False):
            df_dc = df_dc_new.copy()
            c1, c2, c3 = st.columns(3)
            
            with c1:
                if st.session_state.role == "ADMIN":
                    s6_z = st.multiselect("Zone", sorted([x for x in df_dc['ZONE'].unique() if pd.notna(x) and x!=""]), key='z6')
                    if s6_z: df_dc = df_dc[df_dc['ZONE'].isin(s6_z)]
                
                # 🎯 DEPENDENT FILTER
                tu_opts = sorted([x for x in df_dc['TB Unit'].unique() if pd.notna(x) and x!=""])
                s6_tu = st.multiselect("TB Unit", tu_opts, key='tu6')
                if s6_tu: df_dc = df_dc[df_dc['TB Unit'].isin(s6_tu)]
                
                # 🎯 DEPENDENT FILTER
                phi_opts = sorted([x for x in df_dc['PHI'].unique() if pd.notna(x) and x!=""])
                s6_phi = st.multiselect("PHI", phi_opts, key='phi6')
                if s6_phi: df_dc = df_dc[df_dc['PHI'].isin(s6_phi)]
                
            with c2:
                s6_hf = st.multiselect("Facility Type", sorted([x for x in df_dc['Facility_Type'].unique() if pd.notna(x) and x!=""]), key='hf6')
                if s6_hf: df_dc = df_dc[df_dc['Facility_Type'].isin(s6_hf)]
                
                s6_case = st.multiselect("Type of Case", sorted([x for x in df_dc['Type_of_Case'].unique() if pd.notna(x) and x!=""]), key='case6')
                if s6_case: df_dc = df_dc[df_dc['Type_of_Case'].isin(s6_case)]
                
                s6_site = st.multiselect("Site of TBDisease", sorted([x for x in df_dc['Site_of_TBDisease'].unique() if pd.notna(x) and x!=""]), key='site6')
                if s6_site: df_dc = df_dc[df_dc['Site_of_TBDisease'].isin(s6_site)]
                
                s6_outcol = st.multiselect("Treatment Outcome", sorted([x for x in df_dc['Treatment_Outcome'].unique() if pd.notna(x) and x!=""]), key='outcol6')
                if s6_outcol: df_dc = df_dc[df_dc['Treatment_Outcome'].isin(s6_outcol)]

            with c3:
                diag_dt6 = st.date_input("Diagnosis Date Range", value=[], key="d1_6")
                init_dt6 = st.date_input("Initiation Date Range", value=[], key="d2_6")
                out_dt6 = st.date_input("Outcome Date Range", value=[], key="d3_6")
                
        if len(diag_dt6) == 2: df_dc = df_dc[pd.to_datetime(df_dc.get('Diagnosis Date'), errors='coerce').notna() & pd.to_datetime(df_dc.get('Diagnosis Date'), errors='coerce').dt.date.between(diag_dt6[0], diag_dt6[1])]
        if len(init_dt6) == 2: df_dc = df_dc[pd.to_datetime(df_dc.get('Initiation Date'), errors='coerce').notna() & pd.to_datetime(df_dc.get('Initiation Date'), errors='coerce').dt.date.between(init_dt6[0], init_dt6[1])]
        if len(out_dt6) == 2: df_dc = df_dc[pd.to_datetime(df_dc.get('Outcome Date'), errors='coerce').notna() & pd.to_datetime(df_dc.get('Outcome Date'), errors='coerce').dt.date.between(out_dt6[0], out_dt6[1])]

        st.markdown("<hr>", unsafe_allow_html=True)
        
        periods_map = {
            'BASELINE': ('BASELINE', 'Elig_BASELINE'),
            '1ST MONTH': ('1ST MONTH|1 MONTH', 'Elig_1ST_MONTH'),
            '2ND MONTH': ('2ND MONTH|2 MONTH', 'Elig_2ND_MONTH'),
            '3RD MONTH': ('3RD MONTH|3 MONTH', 'Elig_3RD_MONTH'),
            '4TH MONTH': ('4TH MONTH|4 MONTH', 'Elig_4TH_MONTH'),
            '5TH MONTH': ('5TH MONTH|5 MONTH', 'Elig_5TH_MONTH'),
            '6TH MONTH': ('6TH MONTH|6 MONTH', 'Elig_6TH_MONTH')
        }
        
        sel_period = st.radio("📌 Select Follow-up Period to View:", list(periods_map.keys()), horizontal=True)
        p_regex, elig_col = periods_map[sel_period]
        g_col = 'TB Unit' if st.session_state.role == "ZONE" or (st.session_state.role == "ADMIN" and 's6_z' in locals() and len(s6_z) > 0) else 'ZONE' if st.session_state.role == "ADMIN" else 'PHI'
        
        def get_dynamic_summary(df, group_col):
            if df.empty: return pd.DataFrame()
            grp = df.groupby(group_col)
            total_pts = grp.size()
            is_elig = df[elig_col].fillna('').astype(str).str.upper().str.contains("ELIG") & ~df[elig_col].fillna('').astype(str).str.upper().str.contains("NOT")
            eligible_pts = df[is_elig].groupby(group_col).size()
            due = df['Due_Status'].fillna('').astype(str).str.upper()
            not_comp = ~due.str.contains("COMPLETED", na=False)
            is_pending = is_elig & not_comp & due.str.contains(p_regex, na=False)
            pending_pts = df[is_pending].groupby(group_col).size()
            
            summary = pd.DataFrame({'Total Patient': total_pts, 'Eligible': eligible_pts, 'Pending': pending_pts}).fillna(0).astype(int)
            summary['Completed'] = summary['Eligible'] - summary['Pending']
            
            total_patient = summary['Total Patient'].sum()
            total_eligible = summary['Eligible'].sum()
            total_completed = summary['Completed'].sum()
            total_pending = summary['Pending'].sum()
            total_pct = (total_completed / total_eligible * 100) if total_eligible > 0 else 0
            
            summary['% Completed'] = ((summary['Completed'] / summary['Eligible']) * 100).fillna(0).round(1)
            summary = summary.reset_index()
            
            main_zones = ['CENTRAL', 'EAST', 'NORTH', 'NORTH WEST', 'SOUTH', 'SOUTH WEST', 'WEST']
            summary['sort_key'] = summary[group_col].apply(lambda x: main_zones.index(x) if x in main_zones else 998 if x == 'MAPPING NOT DONE' else 999)
            summary = summary.sort_values('sort_key').drop(columns=['sort_key'])
            
            total_row = pd.DataFrame({group_col: ['AMC TOTAL'], 'Total Patient': [total_patient], 'Eligible': [total_eligible], 'Completed': [total_completed], 'Pending': [total_pending], '% Completed': [round(total_pct, 1)]})
            return pd.concat([summary, total_row], ignore_index=True)

        summary_df = get_dynamic_summary(df_dc, g_col)
        
        main_zones = ['CENTRAL', 'EAST', 'NORTH', 'NORTH WEST', 'SOUTH', 'SOUTH WEST', 'WEST']
        
        # 🎯 7 MINI BOXES
        if st.session_state.role == "ADMIN" and ('s6_z' not in locals() or len(s6_z) == 0):
            st.markdown(f"##### 🎯 {sel_period} - Zone Wise % Completed")
            cols7 = st.columns(7)
            for i, z in enumerate(main_zones):
                z_row = summary_df[summary_df[g_col] == z]
                pct_val = 0
                if not z_row.empty: pct_val = z_row['% Completed'].values[0]
                
                if pct_val >= 75: bg_c, t_c = "#d4edda", "#155724" # Green
                elif pct_val >= 50: bg_c, t_c = "#fff3cd", "#856404" # Yellow
                else: bg_c, t_c = "#f8d7da", "#721c24" # Red
                
                card_html = f"""<div style="background-color: {bg_c}; color: {t_c}; border-radius: 5px; padding: 6px 1px; margin-bottom: 10px; text-align: center; border: 1px solid rgba(0,0,0,0.1);"><div style="font-size: 10px; font-weight: bold; text-transform: uppercase;">{z}</div><div style="font-size: 16px; font-weight: 900; margin-top: 2px;">{pct_val}%</div></div>"""
                with cols7[i]: st.markdown(card_html, unsafe_allow_html=True)

        st.markdown(f"##### 📊 {sel_period} Summary ({g_col} Wise)")

        # TARGETED TABLE COLORING
        def color_table(df):
            style_df = pd.DataFrame('', index=df.index, columns=df.columns)
            for i in df.index:
                zone_val = df.at[i, g_col]
                if zone_val in main_zones:
                    try:
                        val_str = str(df.at[i, '% Completed']).replace('%', '')
                        val = float(val_str)
                        if val >= 75:
                            style_df.at[i, '% Completed'] = 'background-color: #d4edda; color: #155724; font-weight: bold;'
                        elif val >= 50:
                            style_df.at[i, '% Completed'] = 'background-color: #fff3cd; color: #856404; font-weight: bold;'
                        else:
                            style_df.at[i, '% Completed'] = 'background-color: #f8d7da; color: #721c24; font-weight: bold;'
                    except: pass
            return style_df

        sum_disp = summary_df.copy()
        sum_disp['% Completed'] = sum_disp['% Completed'].astype(str) + '%'
        
        styled_df = sum_disp.style.apply(color_table, axis=None)
        st.dataframe(styled_df, use_container_width=True, hide_index=True)

        st.markdown(f"##### 📋 {sel_period} Pending Line List")
        is_elig_ll = df_dc[elig_col].fillna('').astype(str).str.upper().str.contains("ELIG") & ~df_dc[elig_col].fillna('').astype(str).str.upper().str.contains("NOT")
        due_ll = df_dc['Due_Status'].fillna('').astype(str).str.upper()
        not_comp_ll = ~due_ll.str.contains("COMPLETED", na=False)
        is_pending_ll = is_elig_ll & not_comp_ll & due_ll.str.contains(p_regex, na=False)
        
        df_ll = df_dc[is_pending_ll].copy()
        if not df_ll.empty:
            ll_cols = ['ZONE', 'TB Unit', 'PHI', 'Type_of_Case', 'Episode ID', 'Patient Name', 'Diagnosis Date', 'Initiation Date', 'Outcome Date', 'Treatment_Outcome', 'Due_Status']
            df_ll_display = df_ll[ll_cols].rename(columns={'Type_of_Case': 'Patient Type', 'Treatment_Outcome': 'Outcome', 'Due_Status': 'Pending Status'})
            st.dataframe(df_ll_display, use_container_width=True, hide_index=True)
            st.download_button(f"📥 Download {sel_period} Pending List", convert_df_to_excel(df_ll_display, f"{sel_period}_Pending"), f"DiffCare_{sel_period}_Pending.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", key=f'dl_ll_{sel_period}')
        else:
            st.success(f"🎉 No pending patients for {sel_period} in the selected criteria!")

        # -------------------------------------------------------------
        # 🎯 NEW ADDITION (MIDDLE): DYNAMIC COHORT MATRIX (UPGRADED FILTERS + TOGGLE)
        # -------------------------------------------------------------
        import datetime
        from dateutil.relativedelta import relativedelta
        
        st.markdown("<br><hr style='border: 1.5px solid #2C3E50;'>", unsafe_allow_html=True)
        st.markdown("<h4 style='color: #2C3E50;'>📊 Consolidated Monthly Pending Matrix (Dynamic Cohorts)</h4>", unsafe_allow_html=True)
        st.markdown("<div style='font-size: 13px; color: #555; margin-bottom: 10px;'><i>Calculates pending patients dynamically based on their specific diagnosis month relative to the review month.</i></div>", unsafe_allow_html=True)
        
        # 🎯 UPGRADED UI: Group By Toggle + All New Filters neatly arranged
        cm1, cm2, cm3 = st.columns(3)
        with cm1:
            mat_view = st.selectbox("📊 View Matrix By", ["Zone", "TB Unit"], key="mat_view_mid", help="Switch between seeing rows by 7 Zones or all 23 TB Units")
            mat_fac = st.selectbox("🏥 Facility Type", ["Public", "Private", "All"], key="mat_fac_mid")
        with cm2:
            today_date = datetime.date.today()
            ref_date = st.date_input("📅 Select Current Review Month", value=today_date, key="mat_ref_dt")
            case_opts = sorted([x for x in df_dc_new['Type_of_Case'].unique() if pd.notna(x) and x!=""])
            mat_case = st.multiselect("Type of Case", case_opts, key="mat_case_mid")
        with cm3:
            tu_opts_mid = sorted([x for x in df_dc_new['TB Unit'].unique() if pd.notna(x) and x!=""])
            mat_tu = st.multiselect("Filter TB Unit", tu_opts_mid, key="mat_tu_mid")
            site_opts = sorted([x for x in df_dc_new['Site_of_TBDisease'].unique() if pd.notna(x) and x!=""])
            mat_site = st.multiselect("Site of TBDisease", site_opts, key="mat_site_mid")

        df_mat = df_dc_new.copy()
        
        if mat_fac == "Public":
            df_mat = df_mat[df_mat['Facility_Type'].astype(str).str.upper().isin(['PUBLIC', 'PHI'])]
        elif mat_fac == "Private":
            df_mat = df_mat[df_mat['Facility_Type'].astype(str).str.upper().isin(['PRIVATE'])]

        # 🎯 APPLY THE NEW MULTISELECT FILTERS
        if mat_tu:
            df_mat = df_mat[df_mat['TB Unit'].isin(mat_tu)]
        if mat_case:
            df_mat = df_mat[df_mat['Type_of_Case'].isin(mat_case)]
        if mat_site:
            df_mat = df_mat[df_mat['Site_of_TBDisease'].isin(mat_site)]

        mat_periods = [
            ('Baseline', 'BASELINE', 'Elig_BASELINE', 1),
            ('1 MONTH', '1ST MONTH|1 MONTH', 'Elig_1ST_MONTH', 2),
            ('2 MONTH', '2ND MONTH|2 MONTH', 'Elig_2ND_MONTH', 3),
            ('3 MONTH', '3RD MONTH|3 MONTH', 'Elig_3RD_MONTH', 4),
            ('4 MONTH', '4TH MONTH|4 MONTH', 'Elig_4TH_MONTH', 5),
            ('5 MONTH', '5TH MONTH|5 MONTH', 'Elig_5TH_MONTH', 6),
            ('6 MONTH', '6TH MONTH|6 MONTH', 'Elig_6TH_MONTH', 7)
        ]
        
        # 🎯 DYNAMIC ROW GROUPING (Zone vs TB Unit)
        if mat_view == "Zone":
            display_entities = ['SOUTH', 'NORTH', 'EAST', 'WEST', 'CENTRAL', 'NORTH WEST', 'SOUTH WEST']
            entity_label = 'ZONE'
            
            def get_zone_mat(z):
                raw_z = str(z).upper().replace("ZONE", "").strip()
                if "SOUTH WEST" in raw_z: return "SOUTH WEST"
                if "NORTH WEST" in raw_z: return "NORTH WEST"
                if "WEST" in raw_z: return "WEST"
                if "SOUTH" in raw_z: return "SOUTH"
                if "CENTRAL" in raw_z: return "CENTRAL"
                if "EAST" in raw_z: return "EAST"
                if "NORTH" in raw_z: return "NORTH"
                return "AMC"
            
            df_mat['Entity_Col'] = df_mat['ZONE'].apply(get_zone_mat)
        else:
            # Display all TB Units
            display_entities = sorted([x for x in df_mat['TB Unit'].unique() if pd.notna(x) and str(x).strip() != ""])
            entity_label = 'TB Unit'
            df_mat['Entity_Col'] = df_mat['TB Unit'].astype(str).str.strip().str.upper()
        
        mat_rows = []
        for entity in display_entities:
            entity_df = df_mat[df_mat['Entity_Col'] == entity]
            row = {entity_label: entity}
            
            for label, rx, elig_col, m_offset in mat_periods:
                target_start = ref_date.replace(day=1) - relativedelta(months=m_offset)
                target_end = (target_start + relativedelta(months=1)) - datetime.timedelta(days=1)
                
                cohort = entity_df[pd.to_datetime(entity_df.get('Diagnosis Date'), errors='coerce').dt.date.between(target_start, target_end)]
                
                is_elig = cohort[elig_col].fillna('').astype(str).str.upper().str.contains("ELIG") & ~cohort[elig_col].fillna('').astype(str).str.upper().str.contains("NOT")
                elig_cnt = is_elig.sum()
                
                due = cohort['Due_Status'].fillna('').astype(str).str.upper()
                not_comp = ~due.str.contains("COMPLETED", na=False)
                is_pending = is_elig & not_comp & due.str.contains(rx, na=False)
                pend_cnt = is_pending.sum()
                
                pct = round((pend_cnt/elig_cnt*100) if elig_cnt>0 else 0)
                
                row[f'Ep_{label}'] = elig_cnt
                row[f'{label}'] = pend_cnt
                row[f'% {label}'] = f"{pct}%"
            mat_rows.append(row)
            
        # AMC Row (Always displays at the bottom)
        amc_row = {entity_label: 'AMC TOTAL'}
        for label, rx, elig_col, m_offset in mat_periods:
            target_start = ref_date.replace(day=1) - relativedelta(months=m_offset)
            target_end = (target_start + relativedelta(months=1)) - datetime.timedelta(days=1)
            
            amc_cohort = df_mat[pd.to_datetime(df_mat.get('Diagnosis Date'), errors='coerce').dt.date.between(target_start, target_end)]
            
            is_elig = amc_cohort[elig_col].fillna('').astype(str).str.upper().str.contains("ELIG") & ~amc_cohort[elig_col].fillna('').astype(str).str.upper().str.contains("NOT")
            elig_cnt = is_elig.sum()
            
            due = amc_cohort['Due_Status'].fillna('').astype(str).str.upper()
            not_comp = ~due.str.contains("COMPLETED", na=False)
            is_pending = is_elig & not_comp & due.str.contains(rx, na=False)
            pend_cnt = is_pending.sum()
            
            pct = round((pend_cnt/elig_cnt*100) if elig_cnt>0 else 0)
            
            amc_row[f'Ep_{label}'] = elig_cnt
            amc_row[f'{label}'] = pend_cnt
            amc_row[f'% {label}'] = f"{pct}%"
            
        mat_rows.append(amc_row)
        
        df_matrix_final = pd.DataFrame(mat_rows)
        
        rename_dict = {}
        for i, (label, _, _, _) in enumerate(mat_periods):
            rename_dict[f'Ep_{label}'] = "Episode ID" + (" " * i)
        df_matrix_final = df_matrix_final.rename(columns=rename_dict)
        
        def style_matrix(styler):
            styler.set_properties(**{'text-align': 'center'})
            styler.set_properties(subset=[entity_label], **{'text-align': 'left', 'font-weight': 'bold', 'background-color': '#f8f9fa'})
            
            for label, _, _, _ in mat_periods:
                col_name = f'% {label}'
                def color_rule(val):
                    try:
                        v = float(val.replace('%', '').strip())
                        if v >= 10: return 'background-color: #F1948A; color: #721c24; font-weight: bold;' # Red
                        elif v >= 5: return 'background-color: #F9E79F; color: #856404; font-weight: bold;' # Yellow
                        else: return 'background-color: #ABEBC6; color: #155724; font-weight: bold;' # Green
                    except:
                        return ''
                styler.map(color_rule, subset=[col_name])
            return styler
        
        st.dataframe(df_matrix_final.style.pipe(style_matrix), use_container_width=True, hide_index=True)

        # -------------------------------------------------------------
        # 🎯 🔄 DIFF CARE COMPARISON ENGINE (BOTTOM)
        # -------------------------------------------------------------
        st.markdown("<br><hr>", unsafe_allow_html=True)
        st.markdown("<h4 style='color: #E67E22;'>🔄 Diff Care Comparison Engine (Old vs New Sheet)</h4>", unsafe_allow_html=True)
        
        cc1, cc2, cc3 = st.columns(3)
        df_dc_comp_new = df_dc_new.copy()
        df_dc_comp_old = df_dc_old.copy()
        
        with cc1:
            comp_zones = st.multiselect("Filter Zone (Comparison)", sorted([x for x in df_dc_comp_new['ZONE'].unique() if pd.notna(x) and x!=""]), key='dc_comp_zone')
            if comp_zones:
                df_dc_comp_new = df_dc_comp_new[df_dc_comp_new['ZONE'].isin(comp_zones)]
                df_dc_comp_old = df_dc_comp_old[df_dc_comp_old['ZONE'].isin(comp_zones)]
                
        with cc2:
            # 🎯 DEPENDENT FILTER
            tu_opts_comp = sorted([x for x in df_dc_comp_new['TB Unit'].unique() if pd.notna(x) and x!=""])
            comp_tus = st.multiselect("Filter TB Unit (Comparison)", tu_opts_comp, key='dc_comp_tu')
            if comp_tus:
                df_dc_comp_new = df_dc_comp_new[df_dc_comp_new['TB Unit'].isin(comp_tus)]
                df_dc_comp_old = df_dc_comp_old[df_dc_comp_old['TB Unit'].isin(comp_tus)]
                
        with cc3:
            comp_dates = st.date_input("Select Diagnosis Date Range", value=[], key="dc_comp_dates")
            
        run_comp = st.button("🚀 Generate Comparison Matrix", use_container_width=True)

        def parse_comp_date(dt_series):
            return pd.to_datetime(dt_series, format='%d-%m-%Y', errors='coerce').combine_first(pd.to_datetime(dt_series, errors='coerce'))

        if run_comp:
            if len(comp_dates) != 2:
                st.error("⚠️ Please select a valid Start and End Date for comparison.")
            else:
                with st.spinner("Analyzing Old and New Diff Care Sheets..."):
                    s_ts = pd.Timestamp(comp_dates[0])
                    e_ts = pd.Timestamp(comp_dates[1])
                    
                    new_dates = parse_comp_date(df_dc_comp_new.get('Diagnosis Date'))
                    old_dates = parse_comp_date(df_dc_comp_old.get('Diagnosis Date'))
                    
                    df_dc_comp_new = df_dc_comp_new[new_dates.notna() & new_dates.dt.date.between(s_ts.date(), e_ts.date())]
                    df_dc_comp_old = df_dc_comp_old[old_dates.notna() & old_dates.dt.date.between(s_ts.date(), e_ts.date())]

                    def get_dc_pend_dict(df):
                        pend = {}
                        if df.empty: return pend
                        for _, r in df.iterrows():
                            eid = str(r['Episode ID']).strip().upper()
                            due = str(r.get('Due_Status', '')).upper()
                            if "COMPLETED" in due:
                                pend[eid] = []
                                continue
                            cur_p = []
                            for p_name, p_reg in periods_map.items():
                                p_rx = p_reg[0] 
                                if re.search(p_rx, due): cur_p.append(p_name)
                            pend[eid] = cur_p
                        return pend

                    old_dict = get_dc_pend_dict(df_dc_comp_old)
                    new_dict = get_dc_pend_dict(df_dc_comp_new)
                    
                    all_comp_ids = set(list(old_dict.keys()) + list(new_dict.keys()))
                    dc_comp_rows = []
                    
                    for eid in all_comp_ids:
                        if eid in ["", "NAN", "NONE"]: continue
                        po = old_dict.get(eid, [])
                        pn = new_dict.get(eid, [])
                        row = {'Episode ID': eid}
                        has_act = False
                        
                        for p_name in list(periods_map.keys()):
                            in_old = p_name in po
                            in_new = p_name in pn
                            
                            if in_old and in_new: row[p_name] = "🟡 PERSISTENT"; has_act = True
                            elif not in_old and in_new: row[p_name] = "🔴 NEW"; has_act = True
                            elif in_old and not in_new: row[p_name] = "🟢 RESOLVED"; has_act = True
                            else: row[p_name] = ""
                            
                        if has_act:
                            r_new = df_dc_comp_new[df_dc_comp_new['Episode ID'] == eid]
                            r_old = df_dc_comp_old[df_dc_comp_old['Episode ID'] == eid]
                            base = r_new.iloc[0] if not r_new.empty else r_old.iloc[0]
                            row['ZONE'] = base.get('ZONE', '')
                            row['TB Unit'] = base.get('TB Unit', '')
                            row['PHI'] = base.get('PHI', '')
                            row['Patient Name'] = base.get('Patient Name', '')
                            row['Facility Type'] = base.get('Facility_Type', '')
                            row['Diagnosis Date'] = base.get('Diagnosis Date', '')
                            dc_comp_rows.append(row)
                            
                    df_final_comp = pd.DataFrame(dc_comp_rows)
                    
                    if not df_final_comp.empty:
                        front = ['ZONE', 'TB Unit', 'PHI', 'Episode ID', 'Patient Name', 'Facility Type', 'Diagnosis Date']
                        other = [c for c in df_final_comp.columns if c not in front]
                        df_final_comp = df_final_comp[front + other]
                        
                        st.success(f"✅ Comparison Generated Successfully for {comp_dates[0].strftime('%d-%b-%Y')} to {comp_dates[1].strftime('%d-%b-%Y')}!")
                        st.dataframe(df_final_comp, use_container_width=True, hide_index=True)
                        st.download_button("📥 Download Comparison Matrix", convert_df_to_excel(df_final_comp, "DC_Comparison"), f"DiffCare_Comparison_{comp_dates[0]}_to_{comp_dates[1]}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", key='dl_dc_comp')
                    else:
                        st.info(f"👍 No differences (🔴 NEW or 🟢 RESOLVED) found between Old and New data for {comp_dates[0].strftime('%d-%b-%Y')} to {comp_dates[1].strftime('%d-%b-%Y')}.")

# ==========================================
# 🟢 TAB 6: STAFF DIRECTORY (HR COMMAND CENTER - MNC ENTERPRISE EDITION)
# ==========================================
with tab6:
    st.markdown("<h3 style='text-align: center; color: #1e293b; font-weight: 800; font-family: system-ui;'>👥 AMC NTEP Staff Directory</h3>", unsafe_allow_html=True)
    
    @st.cache_data(ttl=600)
    def load_staff_directory():
        import re
        base_url = "https://docs.google.com/spreadsheets/d/1uFaHWm7spYKfpe-yrKhe7SC6GafEFM41w45_TnJ1Miw/export?format=csv&gid="
        
        configs = [
            {"name": "MO-SUPERVISOR", "gid": "1725576011", "name_col": "NAME", "zone_col": "ZONE", "tu_col": None},
            {"name": "MO-MEDICAL COLLEGE", "gid": "1072071070", "name_col": "NAME", "zone_col": None, "tu_col": "TU"},
            {"name": "STS", "gid": "1743236661", "name_col": "NAME", "zone_col": "ZONE", "tu_col": "TB UNIT"},
            {"name": "STLS", "gid": "450506055", "name_col": "NAME", "zone_col": "ZONE", "tu_col": "TB UNIT"},
            {"name": "TBHV", "gid": "1273132313", "name_col": "TBHV", "zone_col": "ZONE", "tu_col": "TU"},
            {"name": "LT", "gid": "755154964", "name_col": "NAME", "zone_col": "ZONE", "tu_col": "TU"},
        ]
        
        all_staff = []
        for cfg in configs:
            try:
                df_raw = pd.read_csv(base_url + cfg["gid"])
                
                # 🎯 DYNAMIC HEADER DETECTION
                h_idx = -1
                if str(cfg["name_col"]).upper() in df_raw.columns.astype(str).str.strip().str.upper():
                    df_s = df_raw.copy()
                    df_s.columns = df_s.columns.astype(str).str.strip().str.upper()
                else:
                    for i in range(5):
                        vals = [str(v).upper().strip() for v in df_raw.iloc[i].values]
                        if str(cfg["name_col"]).upper() in vals:
                            h_idx = i; break
                    if h_idx != -1:
                        df_s = df_raw.iloc[h_idx+1:].copy()
                        df_s.columns = df_raw.iloc[h_idx].astype(str).str.strip().str.upper()
                    else:
                        df_s = df_raw.copy()
                        df_s.columns = df_s.columns.astype(str).str.strip().str.upper()

                df_clean = pd.DataFrame()
                
                name_c = str(cfg["name_col"]).upper()
                df_clean['NAME'] = df_s[name_c] if name_c in df_s.columns else ""
                
                zone_c = str(cfg["zone_col"]).upper() if cfg["zone_col"] else ""
                df_clean['RAW_ZONE'] = df_s[zone_c] if zone_c in df_s.columns else ""
                
                tu_c = str(cfg["tu_col"]).upper()
                df_clean['TB_UNIT'] = df_s[tu_c] if tu_c in df_s.columns else "N/A"
                
                df_clean['RAW_ZONE'] = df_clean['RAW_ZONE'].replace(["", "NAN", "NONE", "NaN", pd.NA], None).ffill()
                df_clean['TB_UNIT'] = df_clean['TB_UNIT'].replace(["", "NAN", "NONE", "NaN", pd.NA], None).ffill()
                
                contact_col = next((c for c in df_s.columns if "CONTACT" in c or "CONTECT" in c or "MOBILE" in c), None)
                df_clean['CONTACT NO'] = df_s[contact_col] if contact_col else "N/A"
                
                # STRICT PHI / MEDICAL COLLEGE EXTRACTION
                phi_col = next((c for c in df_s.columns if any(k in c for k in ["PHI", "UHC", "CHC", "FACIL", "INST", "DOT CENTER", "DMC", "MEDICAL COLLEGE"]) and "EMAIL" not in c), None)
                df_clean['PHI/UHC/CHC'] = df_s[phi_col] if phi_col else "N/A"
                df_clean['PHI/UHC/CHC'] = df_clean['PHI/UHC/CHC'].apply(lambda x: "N/A" if "@" in str(x) else x)

                addr_col = next((c for c in df_s.columns if any(k in c for k in ["ADDRESS", "RESIDENCE", "RESIDENCIAL"]) and "EMAIL" not in c), None)
                df_clean['RESIDENCE ADDRESS'] = df_s[addr_col] if addr_col else "N/A"
                df_clean['RESIDENCE ADDRESS'] = df_clean['RESIDENCE ADDRESS'].apply(lambda x: "N/A" if "@" in str(x) else x)
                
                days_col = next((c for c in df_s.columns if any(k in c for k in ["DAY", "JOB LOCATION"]) and "EMAIL" not in c), None)
                df_clean['WORKING_DAYS'] = df_s[days_col] if days_col else "N/A"

                post_col = next((c for c in df_s.columns if "NHM" in c or "GUHP" in c or "NTEP" in c or "AMC" in c), None)
                df_clean['TYPE_OF_POSTING'] = df_s[post_col] if post_col else "N/A"

                dob_col = next((c for c in df_s.columns if "DOB" in c), None)
                df_clean['DOB'] = df_s[dob_col] if dob_col else "N/A"

                pat_col = next((c for c in df_s.columns if "ON TREATMENT" in c and "PUBLIC" in c), None)
                df_clean['ON_TREATMENT'] = df_s[pat_col] if pat_col else "N/A"

                df_clean['SOURCE_SHEET'] = cfg["name"]
                all_staff.append(df_clean)
            except Exception as e:
                pass 
                
        if all_staff:
            final_df = pd.concat(all_staff, ignore_index=True)
            final_df = final_df.dropna(subset=['NAME'])
            
            final_df['NAME'] = final_df['NAME'].astype(str).str.replace('\r', '\n')
            final_df['NAME'] = final_df['NAME'].str.split('\n')
            final_df = final_df.explode('NAME')
            
            final_df['NAME'] = final_df['NAME'].astype(str).str.upper().str.strip()
            final_df['NAME'] = final_df['NAME'].str.replace(r'\s+', ' ', regex=True)
            final_df = final_df[final_df['NAME'] != ""]
            final_df = final_df[~final_df['NAME'].isin(["NAN", "NONE", "N/A"])]
            
            final_df['CONTACT NO'] = final_df['CONTACT NO'].astype(str).str.replace(r'\.0$', '', regex=True).str.replace(r'[^\d]', '', regex=True)
            final_df['CONTACT NO'] = final_df['CONTACT NO'].replace("", "N/A")

            def assign_strict_zone(row):
                raw_z = str(row['RAW_ZONE']).upper().replace("ZONE", "").strip()
                if raw_z not in ["", "NAN", "NONE", "N/A"]:
                    if "SOUTH WEST" in raw_z: return "South West"
                    if "NORTH WEST" in raw_z: return "North West"
                    if "WEST" in raw_z: return "West"
                    if "SOUTH" in raw_z: return "South"
                    if "CENTRAL" in raw_z: return "Central"
                    if "EAST" in raw_z: return "East"
                    if "NORTH" in raw_z: return "North"
                    return raw_z.title()
                
                tu = str(row['TB_UNIT']).upper().strip()
                if any(x in tu for x in ["JODHPUR", "SARKHEJ", "VEJALPUR", "BOPAL", "MAKARBA"]): return "South West"
                if any(x in tu for x in ["SOLA", "GHATLODIA", "CHANDLODIYA", "THALTEJ", "BODAKDEV", "GOTA"]): return "North West"
                if any(x in tu for x in ["VASNA", "PALDI", "SABARMATI", "NAVRANGPURA", "STADIUM", "VADAJ", "RANIP", "CHANDKHEDA"]): return "West"
                if any(x in tu for x in ["DANILIMDA", "VATVA", "MANINAGAR", "ISANPUR", "BEHRAMPURA", "LAMBHA", "PIPLAJ", "NAROL"]): return "South"
                if any(x in tu for x in ["ASARVA", "SHAHPUR", "JAMALPUR", "DARIYAPUR", "CIVIL", "MADHUPURA", "KHANDIA", "DUDHESHWAR"]): return "Central"
                if any(x in tu for x in ["AMRAIWADI", "BHAIPURA", "VASTRAL", "GOMTIPUR", "VIRATNAGAR", "RAMOL", "NIKOL", "ODHAV"]): return "East"
                if any(x in tu for x in ["BAPUNAGAR", "SAIJPUR", "NARODA", "RAKHIAL", "INDIA COLONY", "NOBLENAGAR", "SARDARNAGAR"]): return "North"
                
                return "AMC"
                
            final_df['ZONE'] = final_df.apply(assign_strict_zone, axis=1)
            
            final_df['TB_UNIT'] = final_df['TB_UNIT'].astype(str).str.upper()
            final_df['TB_UNIT'] = final_df['TB_UNIT'].str.replace(r'I/C\s*', '', regex=True)
            final_df['TB_UNIT'] = final_df['TB_UNIT'].str.replace("/", ", ").str.replace("  ", " ").str.title()
            final_df['TB_UNIT'] = final_df['TB_UNIT'].replace(["", "Nan", "None", "N/A"], "N/A")

            final_df['PHI/UHC/CHC'] = final_df['PHI/UHC/CHC'].astype(str).str.upper().replace(["", "NAN", "NONE", "NaN", pd.NA], "N/A").str.title()
            final_df['RESIDENCE ADDRESS'] = final_df['RESIDENCE ADDRESS'].astype(str).str.upper().replace(["", "NAN", "NONE", "NaN", pd.NA], "N/A").str.title()
            final_df['WORKING_DAYS'] = final_df['WORKING_DAYS'].astype(str).str.replace('\n', ', ').str.strip(',').str.title()
            final_df['TYPE_OF_POSTING'] = final_df['TYPE_OF_POSTING'].astype(str).str.upper().replace(["", "NAN", "NONE", "NaN", pd.NA], "N/A").str.upper()
            final_df['DOB'] = final_df['DOB'].astype(str).str.upper().replace(["", "NAN", "NONE", "NaN", pd.NA, "NAT", "00:00:00"], "N/A").str.title()
            final_df['ON_TREATMENT'] = final_df['ON_TREATMENT'].astype(str).replace(r'\.0$', '', regex=True).replace(["NAN", "NONE", "N/A", "nan"], "-")
            
            final_df['DESIGNATION'] = ""
            final_df['FILTER_DESIG'] = ""

            final_df.loc[final_df['SOURCE_SHEET'] == 'MO-SUPERVISOR', 'DESIGNATION'] = "MEDICAL OFFICER SUPERVISOR"
            final_df.loc[final_df['SOURCE_SHEET'] == 'MO-SUPERVISOR', 'FILTER_DESIG'] = "MO-Supervisor"
            final_df.loc[final_df['SOURCE_SHEET'] == 'MO-MEDICAL COLLEGE', 'DESIGNATION'] = "MEDICAL OFFICER"
            final_df.loc[final_df['SOURCE_SHEET'] == 'MO-MEDICAL COLLEGE', 'FILTER_DESIG'] = "Medical Officer"
            final_df.loc[final_df['SOURCE_SHEET'] == 'STLS', 'DESIGNATION'] = "SENIOR TB LABORATORY SUPERVISOR (STLS)"
            final_df.loc[final_df['SOURCE_SHEET'] == 'STLS', 'FILTER_DESIG'] = "STLS"
            final_df.loc[final_df['SOURCE_SHEET'] == 'STS', 'DESIGNATION'] = "SENIOR TREATMENT SUPERVISOR (STS)"
            final_df.loc[final_df['SOURCE_SHEET'] == 'STS', 'FILTER_DESIG'] = "STS"
            final_df.loc[final_df['SOURCE_SHEET'] == 'TBHV', 'DESIGNATION'] = "TB HEALTH VISITOR (TBHV)"
            final_df.loc[final_df['SOURCE_SHEET'] == 'TBHV', 'FILTER_DESIG'] = "TBHV"
            final_df.loc[final_df['SOURCE_SHEET'] == 'LT', 'DESIGNATION'] = "LABORATORY TECHNICIAN (LT)"
            final_df.loc[final_df['SOURCE_SHEET'] == 'LT', 'FILTER_DESIG'] = "LT"

            falguni_mask = final_df['NAME'].str.contains("FALGUNI")
            final_df.loc[falguni_mask, 'ZONE'] = "HEAD OFFICE"
            final_df.loc[falguni_mask, 'TB_UNIT'] = "Arogya Bhavan"

            # 🎯 STRICT NEW HIERARCHY AS REQUESTED
            def assign_hierarchy(sheet_name):
                if sheet_name == "MO-SUPERVISOR": return 1
                if sheet_name == "MO-MEDICAL COLLEGE": return 2
                if sheet_name == "STLS": return 3 
                if sheet_name == "STS": return 4
                if sheet_name == "LT": return 5
                if sheet_name == "TBHV": return 6
                return 99

            final_df['HIERARCHY'] = final_df['SOURCE_SHEET'].apply(assign_hierarchy)
            
            def construct_job_location(row):
                sheet = row['SOURCE_SHEET']
                if sheet == "MO-SUPERVISOR": loc = row['PHI/UHC/CHC']
                elif sheet in ["STLS", "STS", "MO-MEDICAL COLLEGE"]: loc = row['TB_UNIT']
                elif sheet in ["TBHV", "LT"]: 
                    raw_loc = str(row.get('PHI/UHC/CHC', '')).strip()
                    if raw_loc.upper() not in ["N/A", "NAN", "NONE", ""]: loc = raw_loc
                    else: loc = row['TB_UNIT']
                else: loc = "N/A"
                
                days = str(row['WORKING_DAYS']).strip()
                if days.upper() in ["N/A", "NAN", "NONE", ""]:
                    if sheet in ['MO-SUPERVISOR', 'STLS', 'STS', 'MO-MEDICAL COLLEGE', 'LT']:
                        days = "Monday To Saturday"
                    else: days = ""
                
                if days: return f"{loc} [{days}]"
                return loc

            final_df['JOB_LOC_AND_DAYS'] = final_df.apply(construct_job_location, axis=1)
            
            def merge_locations(loc_series):
                locs = []
                for loc_val in loc_series:
                    val = str(loc_val).strip()
                    if val.upper() not in ["N/A", "NAN", "NONE", ""] and val not in locs:
                        locs.append(val)
                return " & ".join(locs) if locs else ""
                
            def get_first_valid(series):
                for val in series:
                    if str(val).upper() not in ["N/A", "NAN", "NONE", ""]: return val
                return ""

            def sum_on_treatment(series):
                total = 0
                has_val = False
                for val in series:
                    try:
                        clean_v = re.sub(r'[^\d.]', '', str(val))
                        if clean_v:
                            total += float(clean_v)
                            has_val = True
                    except: pass
                return str(int(total)) if has_val else "-"
            
            # 🎯 AGGRESSIVE GROUPBY: Fuses rows if NAME and CONTACT NO match! (REPORTS_TO removed from grouping keys)
            final_df = final_df.groupby(['NAME', 'CONTACT NO', 'DESIGNATION', 'FILTER_DESIG', 'SOURCE_SHEET', 'HIERARCHY']).agg({
                'ZONE': lambda x: ' & '.join(sorted(set([z for z in x if z != "N/A"]))),
                'TB_UNIT': lambda x: ' & '.join(sorted(set([t.title() for t in x if str(t).upper() not in ["N/A", "NAN", "NONE"]]))),
                'PHI/UHC/CHC': lambda x: ' & '.join(sorted(set([p.title() for p in x if str(p).upper() not in ["N/A", "NAN", "NONE"]]))),
                'JOB_LOC_AND_DAYS': merge_locations,
                'ON_TREATMENT': sum_on_treatment,
                'TYPE_OF_POSTING': get_first_valid,
                'DOB': get_first_valid,
                'RESIDENCE ADDRESS': get_first_valid
            }).reset_index()
            
            final_df = final_df.sort_values(by=['HIERARCHY', 'ZONE', 'NAME']).reset_index(drop=True)
            
            # 🎯 CALCULATE REPORTS_TO AFTER ZONES ARE MERGED
            mo_sups = final_df[final_df['SOURCE_SHEET'] == "MO-SUPERVISOR"]
            zone_heads = {}
            for _, r in mo_sups.iterrows():
                zones = str(r['ZONE']).split(' & ')
                n = str(r['NAME']).title()
                for z in zones:
                    z = z.strip()
                    if z not in zone_heads: zone_heads[z] = n
                    elif n not in zone_heads[z]: zone_heads[z] += f" & {n}"

            def assign_reporting(row):
                sheet = row['SOURCE_SHEET']
                zones = str(row['ZONE']).split(' & ')
                name = str(row['NAME']).upper()
                
                if "FALGUNI" in name: return "City TB Officer & MO-DTC"
                if sheet == "MO-SUPERVISOR": return "City TB Officer (Dr. S. K. Patel)"
                
                heads = []
                for z in zones:
                    head = zone_heads.get(z.strip(), "Zonal MO-Supervisor")
                    if head not in heads: heads.append(head)
                combo_heads = " & ".join(heads)
                
                if sheet == "MO-MEDICAL COLLEGE": return f"City TB Officer & {combo_heads}"
                return combo_heads

            final_df['REPORTS_TO'] = final_df.apply(assign_reporting, axis=1)

            # Final Visual Polish
            mo_mask = final_df['SOURCE_SHEET'] == 'MO-SUPERVISOR'
            final_df.loc[mo_mask, 'TB_UNIT'] = final_df.loc[mo_mask, 'ZONE'].apply(lambda z: f"All TB Units of {z.title()} Zone" if str(z) not in ["N/A", ""] else "All TB Units")
            final_df.loc[mo_mask, 'PHI/UHC/CHC'] = final_df.loc[mo_mask, 'ZONE'].apply(lambda z: f"All UHC/CHC/Medical College of {z.title()} Zone" if str(z) not in ["N/A", ""] else "All UHC/CHC/Medical College")
            final_df.loc[mo_mask, 'JOB_LOC_AND_DAYS'] = final_df.loc[mo_mask, 'PHI/UHC/CHC'] + " [Monday To Saturday]"
            
            falguni_mask_final = final_df['NAME'].str.contains("FALGUNI")
            final_df.loc[falguni_mask_final, 'JOB_LOC_AND_DAYS'] = "Arogya Bhavan [Monday To Saturday]"
            
            final_df['DOB'] = final_df['DOB'].apply(lambda x: str(x).split(' ')[0] if x != "" else x)
            final_df['DISPLAY_NAME'] = final_df['NAME'].str.title()
            
            return final_df
        return pd.DataFrame()

    with st.spinner("Loading Enterprise HR Data..."):
        df_staff = load_staff_directory()
    
    if df_staff.empty:
        st.warning("⚠️ Staff Directory data could not be loaded. Please check the Google Sheet link and GIDs.")
    else:
        if st.session_state.role == "ZONE":
            target_clean = st.session_state.target.upper().replace("ZONE", "").strip()
            def staff_zone_check(val):
                v_list = [z.strip().replace("ZONE", "").strip() for z in str(val).upper().replace(',', '&').split('&')]
                return target_clean in v_list
            df_staff = df_staff[df_staff['ZONE'].apply(staff_zone_check)]
            
        elif st.session_state.role in ["TB_UNIT", "TB UNIT", "TU"]:
            def staff_tu_check(val):
                v = str(val).upper().strip()
                t_up = st.session_state.target.upper().strip()
                if t_up == "VADAJ" and ("JUNA" in v or "NAVA" in v): return False
                if t_up == "RANIP" and "NEW" in v: return False
                return t_up in v
            df_staff = df_staff[df_staff['TB_UNIT'].apply(staff_tu_check)]
        
        st.markdown("""
        <style>
        div[data-testid="stTextInput"] input { border-radius: 20px; padding: 10px 20px; border: 1px solid #cbd5e1; }
        </style>
        """, unsafe_allow_html=True)
        
        sc1, sc2, sc3, sc4, sc5 = st.columns(5)
        with sc1: search_q = st.text_input("🔍 Search Name, Number...", "")
        
        all_zones_raw = []
        for z_str in df_staff['ZONE']:
            for z in str(z_str).split(' & '): all_zones_raw.append(z.strip())
        zones = ["All Zones"] + sorted(list(set([z for z in all_zones_raw if z not in ["N/A", "NAN", ""]])))
        with sc2: sel_zone = st.selectbox("🏢 Filter Zone", zones)
        
        raw_tus = df_staff[df_staff['ZONE'].str.contains(sel_zone, case=False, na=False)]['TB_UNIT'] if sel_zone != "All Zones" else df_staff['TB_UNIT']
        all_tu_items = set()
        for tu_str in raw_tus.dropna():
            for t in str(tu_str).split('&'):
                cleaned_t = t.strip()
                if cleaned_t and cleaned_t.upper() not in ["N/A", "NAN", "NONE", ""]: all_tu_items.add(cleaned_t)
                    
        with sc3:
            tus = ["All TB Units"] + sorted(list(all_tu_items))
            sel_tu = st.selectbox("🏥 Filter TB Unit", tus)
            
        with sc4:
            raw_phis = df_staff[df_staff['TB_UNIT'].str.contains(sel_tu, case=False, na=False)]['PHI/UHC/CHC'] if sel_tu != "All TB Units" else df_staff['PHI/UHC/CHC']
            all_phi_items = set()
            for phi_str in raw_phis.dropna():
                for p in str(phi_str).split('&'):
                    cleaned_p = p.strip()
                    if cleaned_p and cleaned_p.upper() not in ["N/A", "NAN", "NONE", ""] and "ALL UHC/CHC" not in cleaned_p.upper(): 
                        all_phi_items.add(cleaned_p)
            phis = ["All PHIs"] + sorted(list(all_phi_items))
            sel_phi = st.selectbox("🚑 Filter PHI", phis)

        with sc5:
            desigs = ["All Designations", "MO-Supervisor", "Medical Officer", "STLS", "STS", "LT", "TBHV"]
            sel_desig = st.selectbox("👨‍⚕️ Designation", desigs)
        
        df_display = df_staff.copy()
        if search_q: df_display = df_display[df_display.apply(lambda row: row.astype(str).str.contains(search_q, case=False, na=False).any(), axis=1)]
        if sel_zone != "All Zones": df_display = df_display[df_display['ZONE'].str.contains(sel_zone, case=False, na=False)]
        if sel_tu != "All TB Units": df_display = df_display[df_display['TB_UNIT'].astype(str).str.contains(sel_tu, case=False, na=False)]
        if sel_phi != "All PHIs": df_display = df_display[df_display['PHI/UHC/CHC'].astype(str).str.contains(sel_phi, case=False, na=False)]
        if sel_desig != "All Designations": df_display = df_display[df_display['FILTER_DESIG'] == sel_desig]
        
        st.markdown(f"<div style='color: #64748b; margin-bottom: 20px; font-weight: 600; font-size: 14px;'>Found {len(df_display)} Profiles</div>", unsafe_allow_html=True)
        
        if not df_display.empty:
            display_table = df_display[['ZONE', 'TB_UNIT', 'PHI/UHC/CHC', 'DISPLAY_NAME', 'DESIGNATION', 'ON_TREATMENT', 'TYPE_OF_POSTING', 'DOB', 'CONTACT NO', 'RESIDENCE ADDRESS', 'JOB_LOC_AND_DAYS']].copy()
            display_table = display_table.rename(columns={
                'ZONE': 'Zone',
                'TB_UNIT': 'TB Unit',
                'PHI/UHC/CHC': 'PHI/UHC/CHC',
                'DISPLAY_NAME': 'Name',
                'DESIGNATION': 'Designation',
                'ON_TREATMENT': 'On Treatment (Public)',
                'TYPE_OF_POSTING': 'Type of Posting',
                'DOB': 'DOB',
                'CONTACT NO': 'Mobile No.',
                'RESIDENCE ADDRESS': 'Residence Address',
                'JOB_LOC_AND_DAYS': 'Job Location & Days'
            })
            
            # Final UI Blanks enforcement
            display_table['Mobile No.'] = display_table['Mobile No.'].astype(str).str.replace(r'\.0$', '', regex=True)
            display_table['Mobile No.'] = display_table['Mobile No.'].replace(["N/A", "NAN", "NONE", "nan", ""], "Not Provided")
            display_table['Residence Address'] = display_table['Residence Address'].astype(str).replace(["N/A", "NAN", "NONE", "nan", ""], "Not Provided")
            display_table['PHI/UHC/CHC'] = display_table['PHI/UHC/CHC'].replace(["N/A", "NAN", "NONE", "nan"], "")

            st.dataframe(display_table, use_container_width=True, hide_index=True)
            
            st.download_button(
                label="📥 Download Staff Directory Excel",
                data=convert_df_to_excel(display_table, "Staff_Directory"),
                file_name="AMC_NTEP_Staff_Directory.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key='dl_staff_dir'
            )
        else:
            st.info("No staff profiles found matching the current filters.")

# ==========================================
# 🟢 TAB 7: PRESUMPTIVE TB (NEW)
# ==========================================
with tab7:
    st.markdown("<h3 style='color: #1f618d;'>🔬 Presumptive TB Cases</h3>", unsafe_allow_html=True)
    
    if df_pres_t.empty and df_pres_y.empty:
        st.warning("⚠️ No Presumptive TB data found. Please ensure registers are uploaded to the Colab pipeline.")
    else:
        with st.expander("🔽 Filters & Dates", expanded=True):
            df_p_t = df_pres_t.copy()
            df_p_y = df_pres_y.copy()
            
            # 🎯 FIX: Force blanks to #N/A so it perfectly matches Looker Studio's #N/A row
            df_p_t['ZONE'] = df_p_t['ZONE'].replace(["", "NAN", "NAT", "NONE", "NULL", "<NA>", "N/A"], "#N/A").fillna("#N/A")
            df_p_y['ZONE'] = df_p_y['ZONE'].replace(["", "NAN", "NAT", "NONE", "NULL", "<NA>", "N/A"], "#N/A").fillna("#N/A")
            
            c1, c2, c3 = st.columns(3)
            
            with c1:
                if st.session_state.role == "ADMIN":
                    s7_z = clean_selection(st.multiselect("Zone", get_options_with_counts(df_p_t, 'ZONE', 'tab7'), key='z7'))
                    if s7_z: 
                        df_p_t = df_p_t[df_p_t['ZONE'].isin(s7_z)]
                        df_p_y = df_p_y[df_p_y['ZONE'].isin(s7_z)]
                        
                s7_tu = clean_selection(st.multiselect("TB Unit", get_options_with_counts(df_p_t, 'TB Unit', 'tab7'), key='tu7'))
                if s7_tu: 
                    df_p_t = df_p_t[df_p_t['TB Unit'].isin(s7_tu)]
                    df_p_y = df_p_y[df_p_y['TB Unit'].isin(s7_tu)]
                    
            with c2:
                phi_col = 'Spectrum_Enrolment_PHI' if 'Spectrum_Enrolment_PHI' in df_p_t.columns else 'PHI_Clean'
                if phi_col in df_p_t.columns:
                    s7_phi = clean_selection(st.multiselect("PHI", get_options_with_counts(df_p_t, phi_col, 'tab7'), key='phi7'))
                    if s7_phi: 
                        df_p_t = df_p_t[df_p_t[phi_col].isin(s7_phi)]
                        df_p_y = df_p_y[df_p_y[phi_col].isin(s7_phi)]
                        
                s7_fac = st.multiselect("Facility Type", ["PUBLIC", "PRIVATE"], key='fac7')
                if s7_fac:
                    if 'Facility_Type_Extracted' in df_p_t.columns:
                        df_p_t = df_p_t[df_p_t['Facility_Type_Extracted'].isin(s7_fac)]
                        df_p_y = df_p_y[df_p_y['Facility_Type_Extracted'].isin(s7_fac)]
                        
            with c3:
                if 'Spectrum_Presumptive_Till_Date' in df_p_t.columns:
                    till_dates = st.date_input("Spectrum Presumptive Till Date", value=[], key="d7")
                    if len(till_dates) == 2:
                        df_p_t = df_p_t[pd.to_datetime(df_p_t['Spectrum_Presumptive_Till_Date'], errors='coerce').dt.date.between(till_dates[0], till_dates[1])]
                        df_p_y = df_p_y[pd.to_datetime(df_p_y['Spectrum_Presumptive_Till_Date'], errors='coerce').dt.date.between(till_dates[0], till_dates[1])]

        st.markdown("##### 📊 PTB CASES (Comparison Matrix)")
        
        def get_pres_metrics(df):
            if df.empty: return pd.DataFrame()
            
            # 🎯 FIXED LOGIC: Mirrors Looker Studio Formula -> IF(OR(col="No", col=""), 0, 1)
            blank_no_vals = ['NO', 'NAN', 'NONE', '', 'N/A', '<NA>']
            
            df['Microscopy_Yes'] = ~df.get('Microscopy_Offered', pd.Series(dtype=str)).astype(str).str.strip().str.upper().isin(blank_no_vals)
            df['NAAT_Yes'] = ~df.get('Naat_Offered', pd.Series(dtype=str)).astype(str).str.strip().str.upper().isin(blank_no_vals)
            df['Xray_Yes'] = ~df.get('Xray_Offered', pd.Series(dtype=str)).astype(str).str.strip().str.upper().isin(blank_no_vals)
            
            df['Mic_Naat_Yes'] = df['Microscopy_Yes'] | df['NAAT_Yes']
            
            grp = df.groupby('ZONE', dropna=False).agg(
                Episode_ID=('Episode_ID', 'count'),
                Microscopy=('Microscopy_Yes', 'sum'),
                NAAT=('NAAT_Yes', 'sum'),
                Xray_Offered=('Xray_Yes', 'sum'),
                Mic_Plus_Naat=('Mic_Naat_Yes', 'sum')
            ).reset_index()
            
            grp['%_XRAY_OFFERED'] = (grp['Xray_Offered'] / grp['Episode_ID'] * 100).fillna(0).round(1)
            return grp

        met_t = get_pres_metrics(df_p_t)
        met_y = get_pres_metrics(df_p_y)
        
        if not met_t.empty or not met_y.empty:
            if met_y.empty: met_y = pd.DataFrame(columns=['ZONE', 'Episode_ID', 'Microscopy', 'NAAT', 'Xray_Offered', '%_XRAY_OFFERED', 'Mic_Plus_Naat'])
            if met_t.empty: met_t = pd.DataFrame(columns=['ZONE', 'Episode_ID', 'Microscopy', 'NAAT', 'Xray_Offered', '%_XRAY_OFFERED', 'Mic_Plus_Naat'])
            
            merged = pd.merge(met_y, met_t, on='ZONE', how='outer', suffixes=(' (previous day)', ' (today)')).fillna(0)
            
            req_cols = ['ZONE', 'Episode_ID (previous day)', 'Episode_ID (today)', 'Microscopy (previous day)', 'Microscopy (today)', 
                        'NAAT (previous day)', 'NAAT (today)', 'Xray_Offered (previous day)', 'Xray_Offered (today)', 
                        '%_XRAY_OFFERED (previous day)', '%_XRAY_OFFERED (today)', 'Mic_Plus_Naat (previous day)', 'Mic_Plus_Naat (today)']
            
            for col in req_cols:
                if col not in merged.columns: merged[col] = 0
            
            merged = merged[req_cols]
            
            # Clean Renaming for final display
            rename_dict = {
                'Xray_Offered (previous day)': 'X ray offered yes no (previous day)',
                'Xray_Offered (today)': 'X ray offered yes no (today)',
                '%_XRAY_OFFERED (previous day)': '% XRAY OFFERED (previous day)',
                '%_XRAY_OFFERED (today)': '% XRAY OFFERED (today)',
                'Mic_Plus_Naat (previous day)': 'microscopy + naat (previous day)',
                'Mic_Plus_Naat (today)': 'microscopy + naat (today)'
            }
            merged = merged.rename(columns=rename_dict)
            
            # 🎯 CALCULATE GRAND TOTAL ROW
            numeric_cols = [c for c in merged.columns if 'ZONE' not in c and '%' not in c]
            grand_total = {col: merged[col].sum() for col in numeric_cols}
            grand_total['ZONE'] = 'Grand total'
            
            # Recalculate true percentage for Grand Total
            if grand_total['Episode_ID (previous day)'] > 0:
                grand_total['% XRAY OFFERED (previous day)'] = round((grand_total['X ray offered yes no (previous day)'] / grand_total['Episode_ID (previous day)']) * 100, 1)
            else: grand_total['% XRAY OFFERED (previous day)'] = 0.0
                
            if grand_total['Episode_ID (today)'] > 0:
                grand_total['% XRAY OFFERED (today)'] = round((grand_total['X ray offered yes no (today)'] / grand_total['Episode_ID (today)']) * 100, 1)
            else: grand_total['% XRAY OFFERED (today)'] = 0.0
            
            merged = pd.concat([merged, pd.DataFrame([grand_total])], ignore_index=True)
            
            # Format types
            for col in merged.columns:
                if '%' in col: merged[col] = merged[col].astype(str) + '%'
                elif 'ZONE' not in col: merged[col] = merged[col].astype(int)
            
            # Dynamic Styling specific to XRAY OFFERED colors & Grand Total bolding
            def style_pres_table(styler):
                def color_pct(val):
                    try:
                        v = float(str(val).replace('%', ''))
                        if v < 60: return 'background-color: #E74C3C; color: white;' # Red
                        elif v < 71: return 'background-color: #E67E22; color: white;' # Orange
                        elif v <= 75: return 'background-color: #F1C40F; color: black;' # Yellow
                        else: return 'background-color: #27AE60; color: white;' # Green
                    except: return ''
                
                styler.map(color_pct, subset=['% XRAY OFFERED (previous day)', '% XRAY OFFERED (today)'])
                
                # Make the "Grand total" row bold and visually distinct
                grand_total_idx = merged.index[merged['ZONE'] == 'Grand total']
                if not grand_total_idx.empty:
                    styler.apply(lambda x: ['background-color: #dbeafe; font-weight: bold;' if x.name in grand_total_idx else '' for _ in x], axis=1)
                
                return styler
            
            st.dataframe(merged.style.pipe(style_pres_table), use_container_width=True, hide_index=True)
            
        else:
            st.info("👍 No Presumptive TB records found for the selected filters.")

# ==========================================
# 🟢 TAB 8: ADVERSE OUTCOMES (AUTO-MERGE MASTER TRACKER)
# ==========================================
with tab8:
    st.markdown("<h3 style='color: #0f172a; font-weight: 800; letter-spacing: -0.5px;'>🚨 Auto-Synced Adverse Outcomes Master</h3>", unsafe_allow_html=True)
    st.markdown("<div style='font-size: 13px; color: #555; margin-bottom: 15px;'><i>Automatically detects new adverse outcomes from This Week vs Previous Week, dynamically prevents duplicates, and merges them directly into your Master List!</i></div>", unsafe_allow_html=True)

    @st.cache_data(ttl=600, show_spinner=False)
    def load_adverse_data():
        import urllib.request
        import io
        
        url_master = "https://docs.google.com/spreadsheets/d/1Dfvl87uaZZ12_5F4dhHXTP_u8i9NM9TASWN8wyX18nE/export?format=csv&gid=1027512112"
        url_this = "https://docs.google.com/spreadsheets/d/1Dfvl87uaZZ12_5F4dhHXTP_u8i9NM9TASWN8wyX18nE/export?format=csv&gid=1898426568"
        url_prev = "https://docs.google.com/spreadsheets/d/1Dfvl87uaZZ12_5F4dhHXTP_u8i9NM9TASWN8wyX18nE/export?format=csv&gid=1981365704"

        def get_sheet(url):
            try:
                req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req, timeout=30) as response:
                    df = pd.read_csv(io.BytesIO(response.read()), low_memory=False, dtype=str)
                return df
            except: return pd.DataFrame()

        df_m = get_sheet(url_master)
        if not df_m.empty: df_m.columns = df_m.columns.astype(str).str.strip()

        df_t = get_sheet(url_this)
        if not df_t.empty: df_t.columns = df_t.columns.astype(str).str.strip()
            
        df_p = get_sheet(url_prev)
        if not df_p.empty: df_p.columns = df_p.columns.astype(str).str.strip()

        return df_m, df_t, df_p

    df_master_orig, df_this, df_prev = load_adverse_data()

    # 🛡️ ALIAS MAPPER FOR MASTER SHEET: Maps "REGIME" -> "Type_of_TB_regimen" perfectly
    if not df_master_orig.empty:
        rename_map = {}
        for col in df_master_orig.columns:
            c_up = str(col).strip().upper()
            if c_up == 'AGE': rename_map[col] = 'Age'
            elif c_up in ['REGIME', 'REGIMEN', 'TYPE OF TB REGIMEN', 'TYPE_OF_TB_REGIMEN']: rename_map[col] = 'Type_of_TB_regimen'
        
        df_master_orig = df_master_orig.rename(columns=rename_map)

        # Safety check: Force columns into Master if they are entirely missing
        if 'Age' not in df_master_orig.columns: df_master_orig['Age'] = ""
        if 'Type_of_TB_regimen' not in df_master_orig.columns: df_master_orig['Type_of_TB_regimen'] = ""

    # ---------------------------------------------------------
    # ⚙️ AUTO-MERGE DELTA ENGINE 
    # ---------------------------------------------------------
    df_new_adverse = pd.DataFrame()
    
    import re
    def cx(col_letter):
        num = 0
        for c in col_letter.upper(): num = num * 26 + (ord(c) - ord('A') + 1)
        return num - 1

    def get_col_name_or_idx(df_to_search, possible_names, fallback_col_letter):
        for p in possible_names:
            p_clean = re.sub(r'[^A-Z0-9]', '', str(p).upper())
            for c in df_to_search.columns:
                c_clean = re.sub(r'[^A-Z0-9]', '', str(c).upper())
                if c_clean == p_clean:
                    return c
        idx = cx(fallback_col_letter)
        if idx < len(df_to_search.columns):
            return df_to_search.columns[idx]
        return None

    def safe_extract(df, aliases, col_letter):
        col_name = get_col_name_or_idx(df, aliases, col_letter)
        if col_name in df.columns:
            return df[col_name]
        return pd.Series([""] * len(df), index=df.index)

    if not df_this.empty and not df_prev.empty:
        aliases_id = ['EPISODE ID', 'NTEP ID', 'ID', 'PATIENT ID']
        aliases_out_val = ['TREATMENT OUTCOME', 'OUTCOME']
        aliases_zone = ['ZONE', 'DISTRICT', 'CURRENT DISTRICT', 'CURRENT ZONE', 'SPECTRUM CURRENT ZONE']
        aliases_tu = ['TB UNIT', 'TU', 'CURRENT TU', 'CURRENT TB UNIT', 'SPECTRUM CURRENT TBU']
        aliases_phi = ['SPECTRUM CURRENT HF', 'PHI', 'FACILITY', 'CURRENT PHI', 'HEALTH FACILITY']
        aliases_type = ['FACILITY TYPE', 'TYPE', 'TYPE OF FACILITY', 'SPECTRUM CURRENT HF TYPE']
        aliases_name = ['PATIENT NAME', 'NAME', 'NAME OF PATIENT']
        aliases_diag = ['DIAGNOSIS DATE', 'DATE OF DIAGNOSIS', 'DATE OF TB DIAGNOSIS', 'DX DATE']
        aliases_init = ['INITIATION DATE', 'TREATMENT INITIATION DATE', 'SPECTRUM TREATMENT INITIATION DATE']
        aliases_out = ['OUTCOME DATE', 'DATE OF OUTCOME']

        id_col_this = get_col_name_or_idx(df_this, aliases_id, 'M')
        out_col_this = get_col_name_or_idx(df_this, aliases_out_val, 'BK')
        id_col_prev = get_col_name_or_idx(df_prev, aliases_id, 'M')
        out_col_prev = get_col_name_or_idx(df_prev, aliases_out_val, 'BK')
        
        if id_col_this and out_col_this and id_col_prev and out_col_prev:
            df_this['_ID_UP'] = df_this[id_col_this].fillna("").astype(str).str.strip().str.upper()
            df_this['_OUT_UP'] = df_this[out_col_this].fillna("").astype(str).str.strip().str.upper()
            df_prev['_ID_UP'] = df_prev[id_col_prev].fillna("").astype(str).str.strip().str.upper()
            df_prev['_OUT_UP'] = df_prev[out_col_prev].fillna("").astype(str).str.strip().str.upper()
            
            good = ["CURED", "COMPLETE", "CHANGED", "SUCCESS"]
            blank_variants = ["", "NAN", "N/A", "NONE", "(BLANKS)", "BLANK", "NULL", "nan", "None"]
            
            is_adv_this = ~df_this['_OUT_UP'].str.contains('|'.join(good), na=False)
            has_out_this = ~df_this['_OUT_UP'].isin(blank_variants)
            df_this_adv = df_this[is_adv_this & has_out_this].copy()
            
            prev_map = dict(zip(df_prev['_ID_UP'], df_prev['_OUT_UP']))
            
            def is_new(row):
                pid = row['_ID_UP']
                pout = row['_OUT_UP']
                if pid not in prev_map: return True
                if prev_map[pid] != pout: return True
                return False
                
            df_new = df_this_adv[df_this_adv.apply(is_new, axis=1)].copy()
            df_new = df_new.drop(columns=['_ID_UP', '_OUT_UP'], errors='ignore')
            
            # 🎯 STRICT 14 COLUMNS DEFINITION (Added Age and Regimen)
            master_cols = ['ADVERSE DATE', 'ZONE', 'TB Unit', 'PHI', 'Facility Type', 'Patient Name', 'Episode ID', 'Age', 'Type_of_TB_regimen', 'Diagnosis Date', 'Initiation Date', 'Outcome Date', 'Treatment Outcome', 'On Treatment Days']
            df_export = pd.DataFrame(columns=master_cols)
            
            if not df_new.empty:
                df_export['ZONE'] = safe_extract(df_new, aliases_zone, 'AR')
                df_export['TB Unit'] = safe_extract(df_new, aliases_tu, 'C')
                df_export['PHI'] = safe_extract(df_new, aliases_phi, 'E')
                df_export['Facility Type'] = safe_extract(df_new, aliases_type, 'D')
                df_export['Patient Name'] = safe_extract(df_new, aliases_name, 'N')
                df_export['Episode ID'] = safe_extract(df_new, aliases_id, 'M')
                
                # 🎯 EXACT EXCEL EXTRACTION FOR NEW COLUMNS
                df_export['Age'] = safe_extract(df_new, ['AGE', 'PATIENT AGE'], 'BA')
                df_export['Type_of_TB_regimen'] = safe_extract(df_new, ['TYPE OF TB REGIMEN', 'TB REGIMEN', 'REGIMEN', 'TYPE_OF_TB_REGIMEN', 'REGIME'], 'BJ')
                
                df_export['Diagnosis Date'] = safe_extract(df_new, aliases_diag, 'S')
                df_export['Initiation Date'] = safe_extract(df_new, aliases_init, 'BM')
                df_export['Outcome Date'] = safe_extract(df_new, aliases_out, 'CB')
                df_export['Treatment Outcome'] = safe_extract(df_new, aliases_out_val, 'BK')
                
                # 🛡️ Upgraded Fallback Logic
                def assign_fallback_zone(phi_name, tu_name):
                    val = str(phi_name).upper().strip()
                    if val in ["", "NAN", "NONE", "N/A", "<NA>"]: 
                        val = str(tu_name).upper().strip()
                        if val in ["", "NAN", "NONE", "N/A", "<NA>"]: return ""
                        
                    if any(x in val for x in ["JODHPUR", "SARKHEJ", "VEJALPUR", "BOPAL", "MAKARBA"]): return "SOUTH WEST"
                    if any(x in val for x in ["SOLA", "GHATLODIA", "CHANDLODIYA", "THALTEJ", "BODAKDEV", "GOTA", "TRAGAD", "KD MAIN", "KUSUM"]): return "NORTH WEST"
                    if any(x in val for x in ["VASNA", "PALDI", "SABARMATI", "NAVRANGPURA", "STADIUM", "VADAJ", "RANIP", "CHANDKHEDA", "NHL"]): return "WEST"
                    if any(x in val for x in ["DANILIMDA", "VATVA", "MANINAGAR", "ISANPUR", "BEHRAMPURA", "LAMBHA", "PIPLAJ", "NAROL", "INDRAPURI", "GHODASAR"]): return "SOUTH"
                    if any(x in val for x in ["ASARVA", "SHAHPUR", "JAMALPUR", "DARIYAPUR", "CIVIL", "MADHUPURA", "KHANDIA", "DUDHESHWAR", "KALUPUR"]): return "CENTRAL"
                    if any(x in val for x in ["AMRAIWADI", "BHAIPURA", "VASTRAL", "GOMTIPUR", "VIRATNAGAR", "RAMOL", "NIKOL", "ODHAV", "KHOKHARA"]): return "EAST"
                    if any(x in val for x in ["BAPUNAGAR", "SAIJPUR", "NARODA", "RAKHIAL", "INDIA COLONY", "NOBLENAGAR", "SARDARNAGAR", "MEGHANINAGAR"]): return "NORTH"
                    return "AMC"

                df_export['ZONE'] = df_export.apply(lambda r: assign_fallback_zone(r['PHI'], r['TB Unit']) if str(r['ZONE']).strip().upper() in ["", "NAN", "NONE", "N/A", "<NA>", "NAT"] or len(str(r['ZONE']).strip()) <= 2 else r['ZONE'], axis=1)

                # ⏱️ Multi-Stage Date Parsing Engine
                def clean_date_display(dt_series):
                    s = dt_series.astype(str).str.split(' ').str[0].replace(['nan', 'NaN', 'None', '<NA>', ''], pd.NA)
                    p1 = pd.to_datetime(s, format='%Y-%m-%d', errors='coerce')
                    p2 = pd.to_datetime(s, format='%d-%m-%Y', errors='coerce')
                    p3 = pd.to_datetime(s, dayfirst=True, errors='coerce')
                    
                    valid_dates = p1.combine_first(p2).combine_first(p3)
                    formatted = valid_dates.dt.strftime('%d-%b-%Y')
                    return formatted.fillna(dt_series).replace(['NaT', 'nan', 'NaN', 'None', '<NA>'], '')
                
                df_export['Diagnosis Date'] = clean_date_display(df_export['Diagnosis Date'])
                df_export['Initiation Date'] = clean_date_display(df_export['Initiation Date'])
                df_export['Outcome Date'] = clean_date_display(df_export['Outcome Date'])

                # ⏱️ Accurately calculate On Treatment Days
                today_ts = pd.Timestamp.today(tz='Asia/Kolkata').tz_localize(None).normalize()
                def calc_new_days(row):
                    init = pd.to_datetime(row.get('Initiation Date'), errors='coerce')
                    out = pd.to_datetime(row.get('Outcome Date'), errors='coerce')
                    out_str = str(row.get('Treatment Outcome', '')).upper()
                    if pd.isna(init): return ""
                    return f"{(out - init).days if pd.notna(out) and out_str not in blank_variants else (today_ts - init).days} Days"
                
                df_export['On Treatment Days'] = df_export.apply(calc_new_days, axis=1)
                
                # 🛑 ANTI-DUPLICATE SHIELD
                if not df_master_orig.empty and 'Episode ID' in df_master_orig.columns and 'Treatment Outcome' in df_master_orig.columns:
                    master_keys = df_master_orig['Episode ID'].astype(str).str.strip() + "_" + df_master_orig['Treatment Outcome'].astype(str).str.strip()
                    new_keys = df_export['Episode ID'].astype(str).str.strip() + "_" + df_export['Treatment Outcome'].astype(str).str.strip()
                    df_export = df_export[~new_keys.isin(master_keys)]

                df_new_adverse = df_export.copy()

    # ---------------------------------------------------------
    # 🎛️ UI: DYNAMIC TAGGING (AUTO DAILY DEFAULT) & MASTER MERGING
    # ---------------------------------------------------------
    col_new1, col_new2 = st.columns([1, 2])
    with col_new1:
        st.markdown("<div style='background-color: #f8f9fa; padding: 15px; border-radius: 8px; border-left: 4px solid #b91c1c;'>", unsafe_allow_html=True)
        default_date_tag = pd.Timestamp.today(tz='Asia/Kolkata').strftime('%d %b %Y').upper()
        report_period_input = st.text_input("🏷️ Tag for New Outcomes:", value=default_date_tag, help="This assigns the period tag to the 'ADVERSE DATE' column. It defaults to Today for your daily updates!")
        if not df_new_adverse.empty:
            df_new_adverse['ADVERSE DATE'] = report_period_input
        st.markdown("</div>", unsafe_allow_html=True)
    
    # Merge existing Master with dynamically calculated New records
    df_combined_master = df_master_orig.copy()

    if not df_combined_master.empty and 'Initiation Date' in df_combined_master.columns:
        today_ts_m = pd.Timestamp.today(tz='Asia/Kolkata').tz_localize(None).normalize()
        def calc_master_days(row):
            init = pd.to_datetime(row.get('Initiation Date'), errors='coerce')
            out = pd.to_datetime(row.get('Outcome Date'), errors='coerce')
            out_str = str(row.get('Treatment Outcome', '')).upper()
            if pd.isna(init): return ""
            return f"{(out - init).days if pd.notna(out) and out_str not in ['PENDING', '', 'NAN', 'N/A'] else (today_ts_m - init).days} Days"
        df_combined_master['On Treatment Days'] = df_combined_master.apply(calc_master_days, axis=1)

    if not df_new_adverse.empty:
        df_combined_master = pd.concat([df_combined_master, df_new_adverse], ignore_index=True)
        
    df_combined_master = df_combined_master.replace(["None", "nan", "NaN", "N/A", "<NA>"], "")
    df_combined_master = df_combined_master.fillna("")

    # ---------------------------------------------------------
    # 📊 DASHBOARD FILTERS 
    # ---------------------------------------------------------
    st.write("---")
    c1, c2, c3 = st.columns(3)
    with c1: 
        opts_out = sorted([x for x in df_combined_master['Treatment Outcome'].unique() if str(x).strip() != ""]) if 'Treatment Outcome' in df_combined_master.columns else []
        sel_out = st.multiselect("Filter by Outcome", opts_out, default=opts_out)
    with c2: 
        opts_zone = sorted([x for x in df_combined_master['ZONE'].unique() if str(x).strip() != ""]) if 'ZONE' in df_combined_master.columns else []
        sel_zone = st.multiselect("Filter by Zone", opts_zone)
    with c3: 
        opts_per = sorted([x for x in df_combined_master['ADVERSE DATE'].unique() if str(x).strip() != ""]) if 'ADVERSE DATE' in df_combined_master.columns else []
        sel_period = st.multiselect("Filter by Adverse Date", opts_per)

    df_f = df_combined_master.copy()
    if sel_out and 'Treatment Outcome' in df_f.columns: df_f = df_f[df_f['Treatment Outcome'].isin(sel_out)]
    if sel_zone and 'ZONE' in df_f.columns: df_f = df_f[df_f['ZONE'].isin(sel_zone)]
    if sel_period and 'ADVERSE DATE' in df_f.columns: df_f = df_f[df_f['ADVERSE DATE'].isin(sel_period)]
    
    # ---------------------------------------------------------
    # 📈 METRICS DISPLAY
    # ---------------------------------------------------------
    m1, m2, m3 = st.columns(3)
    with m1:
        st.markdown(f"<div style='background:#fef2f2; border:1px solid #fca5a5; border-radius:10px; padding:15px; text-align:center;'><h4 style='color:#b91c1c; margin:0;'>{len(df_master_orig)}</h4><span style='font-size:12px; color:#7f1d1d;'>Old Master Records</span></div>", unsafe_allow_html=True)
    with m2:
        st.markdown(f"<div style='background:#f0fdf4; border:1px solid #86efac; border-radius:10px; padding:15px; text-align:center;'><h4 style='color:#15803d; margin:0;'>+{len(df_new_adverse)}</h4><span style='font-size:12px; color:#14532d;'>New Auto-Merged Outcomes</span></div>", unsafe_allow_html=True)
    with m3:
        st.markdown(f"<div style='background:#eff6ff; border:1px solid #93c5fd; border-radius:10px; padding:15px; text-align:center;'><h4 style='color:#1d4ed8; margin:0;'>{len(df_f)}</h4><span style='font-size:12px; color:#1e3a8a;'>Total Displayed Records</span></div>", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    
    # 🎯 Ensures final display strictly maps exactly to your requested columns
    display_cols = [c for c in master_cols if c in df_f.columns]
    st.dataframe(df_f[display_cols], use_container_width=True, hide_index=True)
    
    if not df_combined_master.empty:
        st.download_button(
            "📥 Download Auto-Updated Master Excel", 
            convert_df_to_excel(df_combined_master[display_cols], "Master_Adverse"), 
            "Updated_Adverse_Outcomes_Master.xlsx",
            help="Click to download the fully merged file. Simply upload this directly to your Google Drive to replace your old Master Sheet!"
        )

# ==========================================
# 🟢 TAB 9: EPICOLLECT5 LIVE ENTRIES (FIELD DATA)
# ==========================================
with tab9:
    st.markdown("<h3 style='color: #8E44AD;'>📱 Epicollect5 Live Data Sync</h3>", unsafe_allow_html=True)
    st.markdown("<div style='font-size: 13px; color: #555; margin-bottom: 15px;'><i>Live syncing data directly from the Epicollect5 mobile applications.</i></div>", unsafe_allow_html=True)

    # 🎯 Multi-Project Selector (તમારો નવો પ્રોજેક્ટ અહીં Add કર્યો છે)
    projects = {
        "📱 PMDT Patient Visit": "pmdt-patient-visit",
        "🏠 Initial TB Patient Home Visit": "amc-ntep-initial-tb-patient-home-visit-form",
        "🏥 Private Notification Facilities": "private-notification-facilities"
    }
    
    sel_proj_name = st.selectbox("📌 Select Project to View:", list(projects.keys()))
    sel_slug = projects[sel_proj_name]

    @st.cache_data(ttl=300, show_spinner=False)
    def load_epicollect_data(slug):
        import requests
        import io
        import time
        
        all_dfs = []
        page = 1
        error_message = ""
        
        # 🎯 Polite Pagination Engine (Bypasses the 250-entry Anti-Bot Firewall)
        while True:
            url = f"https://five.epicollect.net/api/export/entries/{slug}?format=csv&page={page}"
            try:
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
                response = requests.get(url, headers=headers, timeout=30)
                
                # Handle API Blocks
                if response.status_code == 429:
                    time.sleep(2) # Rate limit hit! Wait 2 seconds and retry this page.
                    continue
                elif response.status_code == 400:
                    if page == 1: error_message = "**HTTP 400:** Epicollect5 rejected the request."
                    break # On later pages, a 400 just means we've exceeded the final page.
                elif response.status_code in [401, 403]:
                    error_message = f"**HTTP {response.status_code}:** Project '{slug}' is PRIVATE. Make it 'Public'."
                    break
                elif response.status_code == 404:
                    error_message = f"**HTTP 404:** Project '{slug}' not found."
                    break
                elif response.status_code != 200:
                    error_message = f"**HTTP {response.status_code}:** API Error."
                    break
                
                # 🎯 Safe CSV Reading (Prevents Pandas crashes on empty final pages)
                try:
                    df = pd.read_csv(io.StringIO(response.text), dtype=str)
                except pd.errors.EmptyDataError:
                    break 
                
                if len(df) == 0: break
                
                cols_to_drop = ['ec5_uuid', 'ec5_parent_uuid', 'ec5_branch_uuid', 'ec5_is_branch', 'title']
                df = df.drop(columns=[c for c in cols_to_drop if c in df.columns], errors='ignore')
                
                all_dfs.append(df)
                
                if len(df) < 50: break # Epicollect's default is 50. If we get 49, it's the last page.
                
                page += 1
                time.sleep(0.5) # ⏳ THE FIX: 0.5s pause to prevent IP ban/firewall drops!
                
            except Exception as e:
                error_message = f"**Connection Error:** {str(e)}"
                break
                
        if not all_dfs: return pd.DataFrame(), error_message
        return pd.concat(all_dfs, ignore_index=True), error_message

    with st.spinner(f"Fetching ALL Live Entries for {sel_proj_name} (Pausing between pages to prevent firewall blocks)..."):
        df_epi_raw, fetch_error = load_epicollect_data(sel_slug)

    if fetch_error and not df_epi_raw.empty:
        st.toast(f"⚠️ Sync paused early: {fetch_error}")

    if df_epi_raw.empty:
        if fetch_error: st.error(fetch_error)
        else: st.warning("⚠️ No data found for this project.")
    else:
        df_epi = df_epi_raw.copy()
        
        col_config = {}
        
        # 1. Clean messy Epicollect column names
        import re
        new_cols = {}
        for c in df_epi.columns:
            clean_name = re.sub(r'^\d+_?', '', c).replace('_', ' ').strip().title()
            new_cols[c] = clean_name
        df_epi = df_epi.rename(columns=new_cols)
        
        # 2. Location to Map Link Converter (MAP LINK FIXED HERE)
        lat_col = next((c for c in df_epi.columns if "Lat" in str(c)), None)
        lon_col = next((c for c in df_epi.columns if "Long" in str(c) or "Lon" in str(c)), None)
        
        if lat_col and lon_col:
            df_epi['📍 Google Map'] = df_epi.apply(
                lambda r: f"https://www.google.com/maps?q={r[lat_col]},{r[lon_col]}" 
                if pd.notna(r[lat_col]) and pd.notna(r[lon_col]) and str(r[lat_col]).strip() != "" else None, 
                axis=1
            )
            col_config['📍 Google Map'] = st.column_config.LinkColumn("📍 Google Map", display_text="🗺️ Open Map")
            df_epi = df_epi.drop(columns=[lat_col, lon_col])

        # 3. Photo Links (Home Visits & PMDT Only)
        if "Home Visit" in sel_proj_name or "PMDT" in sel_proj_name:
            photo_cols = [c for c in df_epi.columns if "Photo" in str(c) or "Image" in str(c) or "Geotagged" in str(c)]
            for p_col in photo_cols:
                df_epi[p_col] = df_epi[p_col].apply(
                    lambda x: f"https://five.epicollect.net/api/export/media/{sel_slug}?type=photo&format=entry_original&name={x}" 
                    if pd.notna(x) and str(x).strip() != "" and not str(x).startswith('http') else x
                )
                col_config[p_col] = st.column_config.ImageColumn(p_col, help="Visit Photo")

        # 4. Filters UI
        st.markdown(f"<div style='background-color:#f4ecf7; padding:15px; border-radius:8px; border: 1px solid #d7bde2; margin-bottom:15px;'><b style='color:#6c3483; font-size:18px;'>Total Submissions: {len(df_epi)}</b></div>", unsafe_allow_html=True)
        
        fc1, fc2, fc3 = st.columns(3)
        with fc1:
            search_epi = st.text_input("🔍 Search Submissions (Name, ID, etc.)", "", key=f"search_epi_{sel_slug}")
        with fc2:
            zone_col = next((c for c in df_epi.columns if "Zone" in str(c)), None)
            if zone_col:
                zones = ["All Zones"] + sorted([str(z) for z in df_epi[zone_col].dropna().unique() if str(z).strip() != ""])
                sel_epi_zone = st.selectbox("🏢 Filter Zone", zones, key=f"zone_{sel_slug}")
            else:
                sel_epi_zone = "All Zones"
        with fc3:
            epi_date_range = st.date_input("📅 Filter by Upload Date Range", value=[], key=f"date_{sel_slug}")

        # 5. Apply Filters
        if search_epi:
            df_epi = df_epi[df_epi.apply(lambda row: row.astype(str).str.contains(search_epi, case=False, na=False).any(), axis=1)]
            
        if sel_epi_zone != "All Zones" and zone_col:
            df_epi = df_epi[df_epi[zone_col].astype(str) == sel_epi_zone]
            
        if len(epi_date_range) == 2:
            upload_col = next((c for c in df_epi.columns if "Upload" in str(c) or "Created" in str(c)), None)
            if upload_col:
                temp_dates = pd.to_datetime(df_epi[upload_col], errors='coerce').dt.date
                df_epi = df_epi[(temp_dates >= epi_date_range[0]) & (temp_dates <= epi_date_range[1])]

        st.markdown(f"<div style='color: #555; margin-bottom: 10px; font-weight: bold;'>Showing {len(df_epi)} Entries</div>", unsafe_allow_html=True)
        
        # 6. Display Table
        st.dataframe(df_epi, column_config=col_config, use_container_width=True, hide_index=True)
        
        st.download_button(
            label=f"📥 Download {sel_proj_name} Data (Excel)",
            data=convert_df_to_excel(df_epi, "Live_Data"),
            file_name=f"{sel_slug}_Live_Data.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f'dl_epi_{sel_slug}'
        )

